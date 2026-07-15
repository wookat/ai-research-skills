"""html_to_pptx.py — convert rendered HTML to native PowerPoint.

Generic pipeline (works on any HTML, not just specific templates):

1. Render in headless chromium → live DOM with computed styles
2. Per element, extract:
   - Tight bounding box of direct text node (via Range API) — separate
     from element's full bounding rect → no parent/child text overlap
   - Computed style: color, font, bg, border, gradient stops, shadow
   - <img> src (data-uri or file://)
3. Convert to PPT shapes:
   - Image      → Picture
   - Decorative (bg / border / gradient / shadow) → Rectangle/RoundedRect
     with matching fill (solid OR linear gradient) + line + shadow
   - Direct text → TextBox tight to text's actual bbox
4. Z-order by DOM depth: containers first (deeper-stacked text on top)

Targets 95%+ visual fidelity. Web-font embedding deferred to Phase 3
(falls back to closest installed font in the meantime).

Usage:
    python -m scripts.html_to_pptx --html poster.html --out poster.pptx
"""
from __future__ import annotations
import argparse
import base64
import json
import re
import sys
import tempfile
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# Native PowerPoint equation support (Insert → Equation editable math).
# Conversion chain: rendered MathJax TeX (data-tex attribute on each
# <mjx-container>) → MathML (latex2mathml) → OMML (mathml2omml) → wrap
# in <a14:m> (DrawingML 2010 math extension; PowerPoint slides only
# recognize OMML inside this wrapper, not as a direct paragraph child).
# Both libs are pip-only pure-Python — fall back silently to the existing
# SVG-rasterization path when they're absent so the script still runs.
try:
    import latex2mathml.converter as _l2m
    import mathml2omml as _m2o
    from lxml import etree as _lxml_etree
    _MATH_AVAILABLE = True
except ImportError:
    _MATH_AVAILABLE = False

# OOXML math namespaces used when injecting native equations.
_M_NS   = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
from lxml import etree

# A0 landscape default (118.9 × 84.1 cm) — change with --width-inch / --height-inch
DEFAULT_W_INCH = 47.0
DEFAULT_H_INCH = 33.1


# CSS generic / Mac-system font names → fonts we install. Both the browser
# (during DOM extract) and PowerPoint (when opening the pptx) often fail to
# resolve these and fall back unpredictably (DejaVu Sans on linux, Calibri on
# Win). Mapping to a known-installed font gives consistent rendering on both
# sides. Browser-side rendering of these aliased fonts requires fontconfig
# alias config too (~/.config/fontconfig/conf.d/99-mac-aliases.conf) so the
# bboxes we measure also use the chosen font.
FONT_ALIASES = {
    # CSS generic + Mac-system + Windows-system font names → fonts likely
    # to be installed (or auto-fontconfig-aliased) on our render machine.
    # font_resolver.py auto-downloads named fonts (Inter, Roboto, etc.) from
    # Google Fonts BEFORE rendering — so for those, we don't need to alias,
    # we use the original family name. The mappings here only cover names
    # that aren't on Google Fonts (Mac/Win OS fonts and CSS keywords).
    "ui-monospace": "JetBrains Mono",
    "ui-sans-serif": "Inter",
    "ui-serif": "Source Serif 4",
    "system-ui": "Inter",
    "-apple-system": "Inter",
    "BlinkMacSystemFont": "Inter",
    "SFMono-Regular": "JetBrains Mono",
    "SF Mono": "JetBrains Mono",
    "Menlo": "JetBrains Mono",
    "Monaco": "JetBrains Mono",
    "Consolas": "JetBrains Mono",
    "Helvetica Neue": "Inter",
    "Helvetica": "Inter",
    "SF Pro Display": "Inter",
    "SF Pro Text": "Inter",
    "Apple Color Emoji": "Inter",
    "Segoe UI": "Inter",
    # NOTE: do NOT alias Cambria / Calibri / Aptos / Arial / Verdana / Georgia
    # / Trebuchet MS / Times New Roman — these are the "POSTER_FONT" cross-
    # platform-safe set. When the user picks one via POSTER_FONT, the walker
    # must preserve the name in PPT XML so PowerPoint (Mac+Win, where these
    # fonts are pre-installed) renders the user's actual choice. Aliasing
    # them would substitute them to whatever happens to be on the Linux
    # render host (e.g. Cambria → Source Serif 4) — wrong on the target.
}


# Font families whose binaries we EMBED into the .pptx (via font_embedder).
# Embedded families render with bit-identical metrics on the browser side
# (where measurement happens) and on the Mac/Win PowerPoint side (where the
# end user opens the file). NON-embedded families measure via host
# fontconfig fallback (Liberation Serif for TNR, DejaVu Sans for Calibri,
# etc.) but render with the host OS's real font in PowerPoint — small
# per-glyph advance differences accumulate, so tight single-line text in a
# non-embedded family needs a small width pad to avoid overflow / unwanted
# wrap on the target.
#
# Update this set when font_embedder.py is taught to embed a new family
# (and the pipeline starts shipping its binaries). Comparison is case-
# insensitive (entries lowercase, the matcher lowercases the run's family
# before comparing).
EMBEDDED_FONT_FAMILIES: set[str] = {"inter"}


def _pick_font_family(font_family_css: str) -> str:
    """Pick the first font from a CSS font-family stack, applying our alias
    map. Walks the list left-to-right until it finds a name we recognize as
    installed (or maps to one), avoiding generic CSS keywords like sans-serif
    that PPT renders as a system default."""
    if not font_family_css:
        return ""
    GENERIC = {"sans-serif", "serif", "monospace", "cursive", "fantasy"}
    for raw in font_family_css.split(","):
        name = raw.strip().strip('"').strip("'")
        if not name:
            continue
        if name in FONT_ALIASES:
            return FONT_ALIASES[name]
        if name.lower() in GENERIC:
            continue
        return name
    return font_family_css.split(",")[0].strip().strip('"').strip("'")


def _split_runs_at_browser_wraps(runs: list[dict], wrap_lines: list[str]) -> list[dict]:
    """Force PPT wrap at exactly the browser's wrap positions.

    Strategy: walk the runs and the wrap_lines in lockstep (matching char-by-
    char, tolerant of whitespace differences). At each line boundary,
    insert a synthetic \\n run so PPT breaks the line there instead of
    relying on its own text shaper to decide where.

    Eliminates wrap drift caused by font-metric differences between
    browser and PPT/soffice (e.g. browser wraps after "1.25%", PPT wraps
    after "1.25" because of subtly different glyph widths).

    Falls back to original runs if shapes don't match (whitespace normalization
    diverges enough to lose alignment) — best-effort, never destructive."""
    if not wrap_lines or len(wrap_lines) <= 1:
        return runs
    # Flat run text (skip existing \n markers)
    run_flat = "".join(r["text"] for r in runs if r["text"] != "\n")
    wrap_flat = "".join(wrap_lines)

    # Whitespace-agnostic equivalence check
    if re.sub(r"\s+", "", run_flat) != re.sub(r"\s+", "", wrap_flat):
        return runs

    # Walk wrap_lines, find char offset in run_flat where each line ends
    cut_offsets = []
    ri = 0
    for line in wrap_lines[:-1]:
        line_chars = line.rstrip()
        if not line_chars:
            continue
        li = 0
        while li < len(line_chars) and ri < len(run_flat):
            r_ch = run_flat[ri]
            l_ch = line_chars[li]
            if r_ch == l_ch:
                ri += 1
                li += 1
            elif r_ch.isspace():
                ri += 1
            elif l_ch.isspace():
                li += 1
            else:
                return runs  # unmatched non-whitespace → abort
        cut_offsets.append(ri)

    if not cut_offsets:
        return runs

    # Walk runs, split + insert \n at each cut
    new_runs = []
    abs_pos = 0
    cut_idx = 0
    # When True, the next non-\n run should have its leading whitespace
    # stripped — set after we insert a synthetic \n right at a run boundary,
    # so the leading space of the *next* run (which would visually appear
    # at the start of the wrapped line in PPT) gets removed. Browsers strip
    # whitespace at wrap points implicitly; PPT does not, so without this
    # the second/third/... lines of wrapped paragraphs render with a small
    # leading indent and the first characters fail to vertically align.
    pending_strip = False
    for r in runs:
        if r["text"] == "\n":
            new_runs.append(r)
            pending_strip = False  # explicit <br>/\n preserves following text
            continue
        text = r["text"]
        local_start = 0
        if pending_strip:
            while local_start < len(text) and text[local_start] in (" ", "\t"):
                local_start += 1
            pending_strip = False
        run_end_abs = abs_pos + len(text)
        # Fire ALL cuts that fall within this run [abs_pos+local_start, run_end_abs]
        while cut_idx < len(cut_offsets) and cut_offsets[cut_idx] <= run_end_abs:
            cut_abs = cut_offsets[cut_idx]
            if cut_abs <= abs_pos + local_start:
                # cut at or before current position — skip (degenerate)
                cut_idx += 1
                continue
            local_cut = cut_abs - abs_pos
            chunk = text[local_start:local_cut]
            if chunk:
                r2 = dict(r); r2["text"] = chunk
                new_runs.append(r2)
            nl = dict(r); nl["text"] = "\n"
            new_runs.append(nl)
            # Skip the whitespace at the wrap point so the wrapped continuation
            # doesn't start with a visible space (this matches browser behavior).
            while local_cut < len(text) and text[local_cut] in (" ", "\t"):
                local_cut += 1
            local_start = local_cut
            cut_idx += 1
        # If we consumed the entire run via cuts ending exactly at the run
        # boundary, the next run's leading whitespace also needs stripping.
        if local_start >= len(text) and new_runs and new_runs[-1].get("text") == "\n":
            pending_strip = True
        if local_start < len(text):
            r3 = dict(r); r3["text"] = text[local_start:]
            new_runs.append(r3)
        abs_pos = run_end_abs
    # Collapse consecutive \n runs. When a synthetic wrap-point inserted
    # above coincides with a pre-existing \n from the source runs (e.g. a
    # source <div class="lbl">MT-Bench<br>3.8B mini</div> where browser
    # also wrapped between the two words because of textbox width), we
    # end up with two adjacent \n runs → PPT renders that as an EMPTY
    # paragraph between them, producing a "mystery line break" the user
    # never asked for. Drop the duplicate. (Bug reported 2026-06-12 on
    # phi-3 stat-mini labels.)
    deduped = []
    prev_was_nl = False
    for r in new_runs:
        if r["text"] == "\n":
            if prev_was_nl:
                continue
            prev_was_nl = True
        else:
            prev_was_nl = False
        deduped.append(r)
    return deduped


def detect_canvas_size(html_path: Path) -> tuple[float, float] | None:
    """Auto-detect physical canvas size in inches from CSS `@page { size: W H }`.

    Many academic posters declare their print size via the CSS @page rule
    (e.g. `@page { size: 60in 36in }` for ACM 60×36 posters, `size: 118.9cm 84.1cm`
    for A0). Parse it before rendering so we can match BOTH the playwright
    viewport AND the PPT slide dimensions to the designer's intent. Without
    this, posters get squashed/stretched into the hardcoded A0 default.

    Returns (width_inch, height_inch) or None if no @page size found.
    Supported units: in, cm, mm, px.
    """
    try:
        text = html_path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return None

    def to_inch(value: str, unit: str) -> float | None:
        try:
            v = float(value)
        except ValueError:
            return None
        u = unit.lower()
        if u == "in":
            return v
        if u == "cm":
            return v / 2.54
        if u == "mm":
            return v / 25.4
        if u == "px":
            return v / 96.0
        if u == "pt":
            return v / 72.0
        return None

    # Match `@page ... { ... size: W<unit> H<unit> ... }` — optionally with
    # named modifiers (landscape/portrait) which we ignore (designer already
    # baked orientation into the W/H order).
    pattern = re.compile(
        r"@page[^{}]*\{[^{}]*?size\s*:\s*"
        r"([0-9.]+)\s*(in|cm|mm|px|pt)\s+([0-9.]+)\s*(in|cm|mm|px|pt)",
        re.IGNORECASE | re.DOTALL,
    )
    m = pattern.search(text)
    if m:
        wi = to_inch(m.group(1), m.group(2))
        hi = to_inch(m.group(3), m.group(4))
        if wi and hi:
            return (wi, hi)
    return None


# OOXML caps each slide dimension at 56 inches (51206400 EMU per axis).
# Posters are commonly 60×36" or 72×48" — bigger than PPT can store as a
# single slide. We cap to PPT_MAX_INCH while preserving aspect ratio.
# SINGLE SOURCE OF TRUTH for the slide cap: both the pipeline
# (auto_correct_loop -> cap_to_pptx_slide) and the standalone main() path
# reference this one constant, so a poster gets the SAME slide size no matter
# which entry point built it. 56 = OOXML hard max (51,206,400 EMU/axis);
# lower it (e.g. 55) for a safety margin and every path follows.
PPT_MAX_INCH = 56.0


def cap_to_pptx_slide(w_inch: float, h_inch: float) -> tuple[float, float]:
    """Scale (w, h) so the longer dimension fits within PPT_MAX_INCH while
    preserving aspect ratio. Used only for SLIDE physical size — viewport
    keeps the original design size for accurate text wrapping."""
    if max(w_inch, h_inch) <= PPT_MAX_INCH:
        return (w_inch, h_inch)
    scale = PPT_MAX_INCH / max(w_inch, h_inch)
    return (w_inch * scale, h_inch * scale)


# ─── color / unit helpers ──────────────────────────────────────────────────

def _parse_color(css: str) -> tuple[RGBColor, float] | None:
    """Parse 'rgb(r,g,b)' / 'rgba(r,g,b,a)' / '#hex' → (RGBColor, alpha).

    Returns None for transparent or unrecognized.

    Alpha compositing: PPT shape lines/fills are flat RGB — there's no
    alpha channel for stroke colors. CSS borders often use rgba with
    alpha < 1 (e.g. `border: 2px solid color-mix(in srgb, var(--accent)
    30%, transparent)` for light tinted underlines). To match the
    rendered visual, we composite the partially-transparent source over
    WHITE (most posters have white/cream backgrounds). The returned
    RGBColor is post-composite; the alpha component is always 1.0
    after this call (caller sees opaque)."""
    if not css or css in ("transparent", "none"):
        return None
    m = re.match(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)(?:\s*,\s*([\d.]+))?\s*\)", css)
    if m:
        r, g, b = (int(x) for x in m.groups()[:3])
        a = float(m.group(4)) if m.group(4) is not None else 1.0
        if a < 0.01:
            return None
        if a < 1.0:
            # Composite over white (#FFFFFF) so the visible color matches
            # what the browser shows on a light background.
            r = int(round(r * a + 255 * (1 - a)))
            g = int(round(g * a + 255 * (1 - a)))
            b = int(round(b * a + 255 * (1 - a)))
        return (RGBColor(r, g, b), 1.0)
    if css.startswith("#"):
        h = css.lstrip("#")
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        if len(h) == 6:
            return (RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16)), 1.0)
    return None


def _px_to_emu(px: float, target_emu: int, viewport_px: float) -> int:
    """Browser px → PPT EMU, scaled proportionally to slide size."""
    return int(px * target_emu / max(1, viewport_px))


def _normalize_color_stop(token: str) -> str | None:
    """Convert a single CSS color token to an rgb()/hex form that
    `_parse_color` understands. Handles:
      - rgb(...) / rgba(...) / #hex — passthrough
      - color(srgb r g b [/ alpha]) — Chromium emits this for color-mix()
        and var() resolution; 0–1 floats per channel → 0–255 ints.
    Returns None when the token is in an unsupported color space.
    """
    s = token.strip()
    if s.startswith("rgb") or s.startswith("#"):
        return s
    # color(srgb 0.893647 0.907294 0.943529 [/ 0.5])
    m = re.match(
        r"^color\(\s*srgb\s+([\d.]+)\s+([\d.]+)\s+([\d.]+)"
        r"(?:\s*/\s*([\d.]+))?\s*\)$",
        s,
    )
    if m:
        r = int(round(float(m.group(1)) * 255))
        g = int(round(float(m.group(2)) * 255))
        b = int(round(float(m.group(3)) * 255))
        a_raw = m.group(4)
        if a_raw is not None:
            a = float(a_raw)
            return f"rgba({r},{g},{b},{a})"
        return f"rgb({r},{g},{b})"
    return None


def _parse_gradient(css: str) -> list[tuple[str, float]] | None:
    """Parse 'linear-gradient(180deg, #f4f7fb 0%, #e8eef6 100%)' →
    [('rgb(244,247,251)', 0.0), ('rgb(232,238,246)', 1.0)].

    Also accepts modern CSS Color 4 stops like
    'color(srgb 0.89 0.90 0.94)' which Chromium emits when the source
    used color-mix() or var() (normalized via _normalize_color_stop).

    Returns None if not a parseable linear gradient."""
    if not css or "linear-gradient" not in css:
        return None
    # Strip 'linear-gradient(' and trailing ')'
    inner = css[css.index("(") + 1: css.rindex(")")]
    # First token is direction (skip if angle/keyword)
    parts = []
    depth = 0
    cur = []
    for c in inner + ",":
        if c == "(":
            depth += 1
        elif c == ")":
            depth -= 1
        if c == "," and depth == 0:
            parts.append("".join(cur).strip())
            cur = []
        else:
            cur.append(c)
    # Drop direction prefix if present (deg / 'to right' etc.)
    if parts and ("deg" in parts[0] or parts[0].startswith("to ") or "turn" in parts[0]):
        parts = parts[1:]
    stops = []
    for i, p in enumerate(parts):
        # Format: '<color> [<pos>%]'. Color may be rgba(...), #hex, or
        # color(srgb r g b [/ a]). Position is optional, defaults evenly.
        m = re.match(
            r"^(rgba?\([^)]+\)|#[0-9a-fA-F]+|color\([^)]+\))"
            r"(?:\s+([\d.]+%))?$",
            p,
        )
        if not m:
            continue
        col_str = _normalize_color_stop(m.group(1))
        if not col_str:
            continue
        pct = m.group(2)
        pct_f = float(pct.rstrip("%")) / 100 if pct else (i / max(1, len(parts) - 1))
        stops.append((col_str, pct_f))
    return stops if len(stops) >= 2 else None


def _set_gradient_fill(shape, stops: list[tuple[str, float]], angle_deg: int = 90) -> None:
    """Manipulate the shape's fill XML to a linear gradient.
    angle_deg: 0 = bottom→top, 90 = left→right, 180 = top→bottom."""
    sp = shape.fill._xPr
    fill_old = sp.find(qn("a:solidFill")) or sp.find(qn("a:noFill")) or sp.find(qn("a:gradFill"))
    if fill_old is not None:
        sp.remove(fill_old)
    grad = etree.SubElement(sp, qn("a:gradFill"), rotWithShape="1", flip="none")
    gs_lst = etree.SubElement(grad, qn("a:gsLst"))
    for col_str, pos in stops:
        c = _parse_color(col_str)
        if not c:
            continue
        rgb, _ = c
        gs = etree.SubElement(gs_lst, qn("a:gs"), pos=str(int(pos * 100000)))
        srgb = etree.SubElement(gs, qn("a:srgbClr"), val=f"{rgb}")
    lin = etree.SubElement(grad, qn("a:lin"), ang=str(angle_deg * 60000), scaled="0")


def _add_shadow(shape, blur_px: float = 5, offset_x_px: float = 0,
                 offset_y_px: float = 2, alpha: float = 0.2) -> None:
    """Add outerShdw to shape's spPr."""
    spPr = shape._element.find(qn("p:spPr")) or shape._element.find(qn("xdr:spPr"))
    if spPr is None:
        # python-pptx textbox or shape; find nested spPr differently
        for child in shape._element.iter():
            if child.tag.endswith("}spPr"):
                spPr = child
                break
    if spPr is None:
        return
    effects = spPr.find(qn("a:effectLst"))
    if effects is None:
        effects = etree.SubElement(spPr, qn("a:effectLst"))
    shdw = etree.SubElement(effects, qn("a:outerShdw"),
                            blurRad=str(int(blur_px * 12700)),
                            dist=str(int((offset_x_px ** 2 + offset_y_px ** 2) ** 0.5 * 9525)),
                            dir=str(int(60000 * (180 if offset_y_px >= 0 else 0))),
                            algn="ctr", rotWithShape="0")
    clr = etree.SubElement(shdw, qn("a:srgbClr"), val="000000")
    alpha_el = etree.SubElement(clr, qn("a:alpha"), val=str(int(alpha * 100000)))


def _set_round_rect_adj(shape, radius_px: float, w_px: float, h_px: float) -> None:
    """PPT roundRect preset uses an `adj` value (fraction of min(w,h) in
    1/100000 units). Default is 16667 = ~16.7% which on a tall box becomes
    huge. To match CSS `border-radius: 4px` we need adj = radius/min(w,h)."""
    min_dim = max(1.0, min(w_px, h_px))
    adj = min(50000, max(0, int(radius_px / min_dim * 100000)))
    sp = shape._element
    for prst in sp.iter(qn("a:prstGeom")):
        avLst = prst.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prst, qn("a:avLst"))
        # Wipe existing gd entries; rewrite ours
        for gd in list(avLst.findall(qn("a:gd"))):
            avLst.remove(gd)
        etree.SubElement(avLst, qn("a:gd"), name="adj", fmla=f"val {adj}")
        break


def _set_run_letter_spacing(run, spacing_px: float, font_size_px: float) -> None:
    """OOXML run rPr `spc` is in 1/100 pt. CSS letter-spacing is px."""
    if abs(spacing_px) < 0.05:
        return
    spc_pt = spacing_px * 72 / 96
    spc_units = int(spc_pt * 100)
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(spc_units))


def _set_run_baseline(run, vertical_align: str) -> None:
    """OOXML run rPr `baseline` is offset percentage * 1000.
    `super` = +30000 (30% raised), `sub` = -25000 (25% lowered).
    Mirrors PowerPoint's built-in Superscript/Subscript text styles."""
    if vertical_align == "super":
        rPr = run._r.get_or_add_rPr()
        rPr.set("baseline", "30000")
    elif vertical_align == "sub":
        rPr = run._r.get_or_add_rPr()
        rPr.set("baseline", "-25000")


def _set_run_highlight(run, color_rgb: RGBColor) -> None:
    """OOXML run rPr `<a:highlight>` — PowerPoints "Text Highlight Color"
    (text highlight color, since Office 2019). PowerPoint validates rPr
    child order strictly: highlight must come BEFORE <a:latin>/<a:ea>/<a:cs>/
    <a:sym> font typeface elements (and AFTER fill elements). Out-of-order
    children are silently dropped by PowerPoint even though soffice tolerates
    them. We insert at the right position rather than just appending."""
    rPr = run._r.get_or_add_rPr()
    # Strip existing highlight (idempotent)
    for child in list(rPr):
        if child.tag == qn("a:highlight"):
            rPr.remove(child)
    hl = etree.Element(qn("a:highlight"))
    etree.SubElement(hl, qn("a:srgbClr"), val=f"{color_rgb}")
    # Per ECMA-376 / DrawingML schema, highlight comes AFTER fill/effectLst
    # and BEFORE <a:uLnTx> ... <a:latin>/<a:ea>/<a:cs>/<a:sym>/<a:hlinkClick>.
    # Insert before the first font-typeface or hyperlink element if present;
    # otherwise just append (no later siblings exist).
    insert_before_tags = {
        qn("a:uLnTx"), qn("a:uLn"), qn("a:uFillTx"), qn("a:uFill"),
        qn("a:latin"), qn("a:ea"), qn("a:cs"), qn("a:sym"),
        qn("a:hlinkClick"), qn("a:hlinkMouseOver"), qn("a:rtl"),
        qn("a:extLst"),
    }
    insert_pos = None
    for i, child in enumerate(rPr):
        if child.tag in insert_before_tags:
            insert_pos = i
            break
    if insert_pos is None:
        rPr.append(hl)
    else:
        rPr.insert(insert_pos, hl)


# ─── DOM extraction ────────────────────────────────────────────────────────

def extract_dom(html_path: Path, viewport_w: int | None = None,
                viewport_h: int | None = None) -> dict:
    """Extract elements + text-block structures.

    Each text-block = one HTML BLOCK element (p, li, h1-h6, figcaption, td, th)
    with its descendants serialized as inline runs (text + style). Becomes one
    PPT TextBox with mixed-style runs in a single paragraph.

    This avoids the previous text-node-per-textbox approach where inline
    children (<strong>, <em>) split the parent text into disconnected boxes
    that didn't visually flow.

    Viewport dimensions: if not provided, auto-detected from the HTML's
    @page size and converted to px @ 96 DPI. Falls back to A0 landscape.
    """
    if viewport_w is None or viewport_h is None:
        detected = detect_canvas_size(html_path)
        if detected:
            wi, hi = detected
            viewport_w = int(wi * 96)
            viewport_h = int(hi * 96)
        else:
            viewport_w = viewport_w or int(DEFAULT_W_INCH * 96)
            viewport_h = viewport_h or int(DEFAULT_H_INCH * 96)
    print(f"      [viewport] {viewport_w}x{viewport_h}px", file=sys.stderr)
    with sync_playwright() as p:
        b = p.chromium.launch()
        # device_scale_factor=2 only affects rasters (the math-region
        # screenshots below); getBoundingClientRect stays in CSS px, so all
        # element/text geometry is unchanged. Gives crisp equation PNGs.
        page = b.new_context(viewport={"width": viewport_w,
                                       "height": viewport_h},
                             device_scale_factor=2).new_page()
        # Serve MathJax from the skill's bundled tex-svg.js instead of the
        # CDN so equation typeset is deterministic on offline/restricted
        # hosts (else formulas silently fail to render and the OMML/SVG
        # extraction below sees no <mjx-container>). Registered BEFORE goto.
        # Falls back to the network if the bundle is absent.
        _mj = (Path(__file__).resolve().parents[2]
               / "assets" / "mathjax" / "tex-svg.js")
        if _mj.is_file():
            def _route_mj(route):
                try:
                    route.fulfill(path=str(_mj),
                                  content_type="application/javascript")
                except Exception:
                    try:
                        route.continue_()
                    except Exception:
                        pass
            try:
                page.route("**/tex-svg.js", _route_mj)
            except Exception:
                pass
        # Same offline mirror for KaTeX (the other math engine): serve css / js /
        # auto-render / woff2 fonts from assets/katex/ so KaTeX-typeset posters
        # render deterministically offline. Matches any katex@<ver>/dist/** URL and
        # maps the path tail into the mirror (contrib/ flattened to top level).
        _katex_dir = (Path(__file__).resolve().parents[2] / "assets" / "katex")
        if _katex_dir.is_dir():
            _ktypes = {".css": "text/css", ".js": "application/javascript",
                       ".woff2": "font/woff2", ".woff": "font/woff", ".ttf": "font/ttf"}
            def _route_katex(route):
                try:
                    m = re.search(r"/dist/(.+)$", route.request.url)
                    tail = (m.group(1) if m else "").split("?")[0].replace("contrib/", "")
                    f = _katex_dir / tail
                    if tail and f.is_file():
                        route.fulfill(path=str(f),
                                      content_type=_ktypes.get(f.suffix, "application/octet-stream"))
                    else:
                        route.continue_()
                except Exception:
                    try:
                        route.continue_()
                    except Exception:
                        pass
            try:
                page.route("**/katex@*/dist/**", _route_katex)
            except Exception:
                pass
        page.goto(f"file://{html_path.resolve()}", wait_until="networkidle")
        page.wait_for_timeout(800)
        # Force every <img> to finish DECODING before we read geometry below.
        # networkidle fires on network quiescence, but an <img> can still be
        # pre-decode: its naturalWidth is 0 and, with width:auto in a flex row
        # (header institution logos, the QR tiles), it collapses to ~0px. That
        # made logos capture at a 16px bbox and dropped the QR <img> entirely via
        # the `r.width < 1` guard below — the "logo misaligned / QR not loading"
        # bugs. Awaiting decode() populates naturalWidth so layout settles.
        try:
            page.evaluate(
                """() => Promise.all([...document.images].map(img => {
                    if (img.complete && img.naturalWidth > 0) return null;
                    if (img.decode) return img.decode().catch(() => {});
                    return new Promise(res => { img.onload = img.onerror = res; });
                }))"""
            )
            page.wait_for_timeout(200)
        except Exception:
            pass
        data = page.evaluate("""() => {
          const elements = [];
          const text_blocks = [];

          // Normalize ANY CSS color string (rgb, oklab, color(), hsl, named,
          // color-mix output, etc.) to "rgba(r, g, b, a)" via a 1x1 canvas.
          // Python only parses rgb/rgba/hex — letting the browser do the
          // resolution avoids needing a full CSS color parser in Python.
          const _toRGB_canvas = document.createElement('canvas');
          _toRGB_canvas.width = _toRGB_canvas.height = 1;
          const _toRGB_ctx = _toRGB_canvas.getContext('2d');
          const toRGB = (cssColor) => {
            if (!cssColor || cssColor === 'transparent' ||
                cssColor === 'none' || cssColor === '') {
              return cssColor || '';
            }
            try {
              _toRGB_ctx.clearRect(0, 0, 1, 1);
              _toRGB_ctx.fillStyle = '#000';
              _toRGB_ctx.fillStyle = cssColor;
              _toRGB_ctx.fillRect(0, 0, 1, 1);
              const d = _toRGB_ctx.getImageData(0, 0, 1, 1).data;
              if (d[3] === 0) return 'transparent';
              return `rgba(${d[0]}, ${d[1]}, ${d[2]}, ${(d[3]/255).toFixed(3)})`;
            } catch (e) { return cssColor; }
          };

          const BLOCK_TAGS = new Set([
            'p','h1','h2','h3','h4','h5','h6','li','figcaption','td','th',
            'blockquote','dt','dd','summary','caption'
          ]);

          // Block-layout-creating tags. A descendant of any of these forces
          // the parent NOT inline-only (since the parent contains stacked
          // block content). CRITICAL: this is tag-name based, NOT CSS-
          // display based. Earlier display-based check broke when an h3 had
          // `display: flex` with span.num (display: block) as the circle
          // marker — span.num was considered "block-like" so h3 dropped its
          // own title text. Section titles disappeared. Using tag names
          // matches semantic intent: a <span> is inline-flow regardless of
          // its display, a <div> is block-flow regardless of its display.
          const BLOCK_LAYOUT_TAGS = new Set([
            'div','p','h1','h2','h3','h4','h5','h6','ul','ol','li',
            'figure','figcaption','section','article','header','footer',
            'nav','main','aside','blockquote','table','thead','tbody',
            'tr','td','th','pre','address','dl','dt','dd','form',
            'fieldset','legend','hr'
          ]);

          // Inline-only descendants? Tag-name based: any block-layout tag
          // among descendants → not inline-only. Span/em/strong/a/etc are
          // inline-flow even if styled as display:block.
          const isInlineOnly = (el) => {
            for (const c of el.querySelectorAll('*')) {
              if (BLOCK_LAYOUT_TAGS.has(c.tagName.toLowerCase())) return false;
            }
            return true;
          };

          // A flex/grid row of COMPOSITE items (each child is an image chip or
          // its own flex/grid box) — e.g. the header .strip of logo chips, a
          // stat-strip of cells. isInlineOnly only looks at descendant TAG names
          // (span/img aren't block tags), so it wrongly reports such a row as an
          // inline-only text line and the whole row gets emitted as ONE centered
          // text block. That is invisible when the row is all images (no text),
          // but when the row's only text is a lone label (the venue-logo
          // fallback "<Venue> <Year>" when no conf logo resolved), the label
          // renders centered across the full row, overlapping the logos. Exclude
          // these rows from text emission; their text-bearing leaves emit
          // themselves (see isFlexLabelSpan's ancestor walk).
          const isCompositeRow = (node) => {
            const d = getComputedStyle(node).display;
            if (!['flex', 'inline-flex', 'grid', 'inline-grid'].includes(d)) return false;
            if (node.childElementCount < 2) return false;
            let boxy = 0;
            for (const c of node.children) {
              const cd = getComputedStyle(c).display;
              if (c.querySelector('img')
                  || ['flex', 'inline-flex', 'grid', 'inline-grid'].includes(cd)) {
                boxy += 1;
              }
            }
            return boxy >= 2;
          };

          // Chip-pill predicate (mirrors the Python detector). Used by
          // collectRuns to skip element-node descendants that will render
          // as their own RoundedRect chip-pill, AND by the tight-bbox
          // recompute below text_blocks.push to detect when a parent's
          // bbox over-covers the surviving text region. Hoisted to outer
          // scope so both call sites share one definition.
          const isChipPillNode = (node) => {
            if (!['span', 'div'].includes(node.tagName.toLowerCase())) return false;
            if (node.children.length > 0) return false;
            const ncs = getComputedStyle(node);
            const radius = parseFloat(ncs.borderTopLeftRadius) || 0;
            if (radius < 4) return false;
            const bg = toRGB(ncs.backgroundColor);
            if (!bg || bg === 'transparent' || bg.includes(', 0)')) return false;
            const txt = node.textContent.trim();
            if (txt.length === 0 || txt.length > 60) return false;
            const pdisp = node.parentElement
              ? getComputedStyle(node.parentElement).display : '';
            if (!['flex', 'inline-flex', 'grid', 'inline-grid'].includes(pdisp)) return false;
            const rr = node.getBoundingClientRect();
            if (rr.height > 150 || rr.width > 600) return false;
            return true;
          };

          // Walk text nodes and measure per-character Y position via Range
          // API to detect WHERE the browser actually wrapped. Returns the
          // text split into lines exactly where wrapping occurred. Used to
          // force PPT to wrap at the same points (we insert newlines),
          // eliminating wrap drift from font-shaper differences between
          // browser and PPT/soffice.
          const extractWrapLines = (el) => {
            const walker = document.createTreeWalker(el, NodeFilter.SHOW_TEXT);
            const textNodes = [];
            let n;
            while (n = walker.nextNode()) {
              if (n.textContent && n.textContent.length > 0) textNodes.push(n);
            }
            if (!textNodes.length) return [];

            const range = document.createRange();
            const lines = [];
            let curLine = "";
            // Threshold = half the element's CSS line-height. Real wraps
            // jump by ~1 line-height; subscript/superscript baseline shifts
            // are typically <30% of line-height. Half-line-height threshold
            // cleanly separates real wraps from sub/sup y-shifts without
            // false splits on math like J<sub>EE</sub>.
            const cs = getComputedStyle(el);
            const lhPx = parseFloat(cs.lineHeight) || (parseFloat(cs.fontSize) * 1.2);
            const LINE_THRESHOLD = Math.max(20, lhPx * 0.5);
            let lineBaseline = null;

            for (const tn of textNodes) {
              const text = tn.textContent;
              const len = text.length;
              if (len === 0) continue;
              for (let i = 0; i < len; i++) {
                range.setStart(tn, i);
                range.setEnd(tn, i + 1);
                const r = range.getBoundingClientRect();
                if (r.width === 0 && r.height === 0) {
                  curLine += text[i];
                  continue;
                }
                const y = Math.round(r.top);
                if (lineBaseline === null) {
                  lineBaseline = y;
                  curLine += text[i];
                } else if (y > lineBaseline + LINE_THRESHOLD) {
                  if (curLine.length > 0) lines.push(curLine);
                  curLine = text[i];
                  lineBaseline = y;
                } else {
                  curLine += text[i];
                }
              }
            }
            if (curLine.length > 0) lines.push(curLine);
            return lines;
          };

          // Serialize an element's descendants as inline runs
          // Each run: {text, font_family, font_size_px, font_weight,
          //            font_style, color, letter_spacing_px,
          //            vertical_align, bg_color}
          //
          // Pseudo-element content via `::before` / `::after { content: "..." }`
          // is emitted as a synthetic run prepended/appended around the
          // children. This catches "So what →" callout prefixes,
          // bullet markers via `content: "▶"`, status icons, etc.
          // Skip if content is `none` / `normal` (browser default), or if
          // the pseudo-element's display is none.
          const _unescapeCssContent = (s) => {
            if (!s) return '';
            // Strip surrounding quotes
            let v = s;
            if ((v.startsWith('"') && v.endsWith('"')) ||
                (v.startsWith("'") && v.endsWith("'"))) {
              v = v.slice(1, -1);
            }
            // Decode \\xxxx CSS unicode escapes → real char
            return v.replace(/\\\\([0-9a-fA-F]{1,6})\\s?/g,
              (_, h) => String.fromCodePoint(parseInt(h, 16)));
          };
          const _pseudoRun = (el, which) => {
            // which = '::before' or '::after'
            const pcs = getComputedStyle(el, which);
            if (!pcs) return null;
            if (pcs.display === 'none' || pcs.visibility === 'hidden') return null;
            const raw = pcs.content;
            if (!raw || raw === 'none' || raw === 'normal') return null;
            // Suppress as a text run when the pseudo has visual chrome
            // (bg + non-trivial border-radius) — it will be emitted as a
            // separate decorated chip-pill shape by _emitPseudoDeco below,
            // so emitting it ALSO as a text run would double-render the
            // content (e.g. .p-steps .step::before circles).
            const _pseudoBg = toRGB(pcs.backgroundColor);
            const _pseudoR = parseFloat(pcs.borderTopLeftRadius) || 0;
            if (_pseudoBg && _pseudoBg !== 'transparent'
                && !_pseudoBg.includes(', 0)') && _pseudoR >= 2) {
              return null;
            }
            const text = _unescapeCssContent(raw);
            if (!text) return null;
            return {
              text,
              font_family: pcs.fontFamily,
              font_size_px: parseFloat(pcs.fontSize),
              font_weight: pcs.fontWeight,
              font_style: pcs.fontStyle,
              color: toRGB(pcs.color),
              letter_spacing_px: pcs.letterSpacing === 'normal' ? 0 : (parseFloat(pcs.letterSpacing) || 0),
              text_transform: pcs.textTransform,
              vertical_align: 'baseline',
              bg_color: '',
            };
          };
          const collectRuns = (el, runs) => {
            // (isChipPillNode is hoisted to outer scope so the tight-bbox
            // recompute below text_blocks.push can reuse the same predicate.)
            const before = _pseudoRun(el, '::before');
            if (before) runs.push(before);
            for (const node of el.childNodes) {
              if (node.nodeType === Node.TEXT_NODE) {
                const text = node.textContent;
                if (!text) continue;
                const pcs = getComputedStyle(el);
                // vertical-align: super/sub (or <sup>/<sub> ancestor) → PPT
                // baseline offset attribute on the run, not a separate font
                // size change. CSS reports computed value as keyword.
                const va = pcs.verticalAlign;
                let vertical_align = 'baseline';
                if (va === 'super' || el.tagName === 'SUP') vertical_align = 'super';
                else if (va === 'sub' || el.tagName === 'SUB') vertical_align = 'sub';
                else {
                  // walk ancestors to find sup/sub (computed va is on the
                  // intermediate element, not necessarily the text-holder)
                  let p = el;
                  while (p && p !== document.body) {
                    const t = p.tagName;
                    if (t === 'SUP') { vertical_align = 'super'; break; }
                    if (t === 'SUB') { vertical_align = 'sub'; break; }
                    p = p.parentElement;
                  }
                }
                // Inline bg color: if THIS element has a non-transparent bg
                // and is inline-displayed, capture as RUN highlight (not as
                // a separate decorated rectangle).
                let bg_color = '';
                const cs_disp = pcs.display;
                if ((cs_disp === 'inline' || cs_disp === 'inline-block')
                    && pcs.backgroundColor) {
                  const bg = toRGB(pcs.backgroundColor);
                  if (bg && bg !== 'transparent' && !bg.includes(', 0)')) {
                    bg_color = bg;
                  }
                }
                runs.push({
                  text,
                  font_family: pcs.fontFamily,
                  font_size_px: parseFloat(pcs.fontSize),
                  font_weight: pcs.fontWeight,
                  font_style: pcs.fontStyle,
                  color: toRGB(pcs.color),
                  letter_spacing_px: pcs.letterSpacing === 'normal' ? 0 : (parseFloat(pcs.letterSpacing) || 0),
                  text_transform: pcs.textTransform,
                  vertical_align,
                  bg_color,
                });
              } else if (node.nodeType === Node.ELEMENT_NODE) {
                const cs = getComputedStyle(node);
                if (cs.display === 'none' || cs.visibility === 'hidden') continue;
                if (node.tagName.toLowerCase() === 'br') {
                  runs.push({text: '\\n', font_family: '', font_size_px: 0,
                             font_weight: '', font_style: '', color: '',
                             letter_spacing_px: 0, text_transform: 'none',
                             vertical_align: 'baseline', bg_color: ''});
                } else if (isChipPillNode(node)) {
                  // Chip-pill descendant: its text is rendered inside the
                  // styled RoundedRect's text_frame by the Python builder.
                  // Don't recurse — otherwise the chip text doubles into
                  // the parent text_block (e.g. P17 .p-banner: "Blocked"
                  // tag would appear both inside the red pill AND as plain
                  // text in the banner's full-width textbox).
                  continue;
                } else if (node.tagName.toLowerCase() === 'mjx-container'
                           || (node.classList && node.classList.contains('katex'))) {
                  // Math container (MathJax <mjx-container> OR KaTeX <span
                  // class="katex">) — the equation is emitted as a rasterized
                  // Picture (Pass 3) at its absolute bbox. Don't recurse into
                  // it: its hidden a11y tree — MathJax's <mjx-assistive-mml> or
                  // KaTeX's .katex-mathml — plus KaTeX's per-glyph .katex-html
                  // spans would leak the equation as garbled text into the
                  // parent textbox, drawing it twice.
                  // BUT for INLINE math we must RESERVE the equation's width in
                  // the surrounding text, else the following text reflows left
                  // as if the equation were zero-width and the Picture overlaps
                  // it (observed: "...h_card" then "FULL band" collide). Push a
                  // non-breaking-space spacer sized to the equation width so the
                  // text leaves a gap the Picture drops into. Display math sits
                  // on its own line (no shared text) so it needs no spacer.
                  const _isDisplayMath = node.tagName.toLowerCase() === 'mjx-container'
                    ? node.getAttribute('data-display') === '1'
                    : (node.closest('.katex-display') !== null);
                  if (!_isDisplayMath) {
                    const _mrect = node.getBoundingClientRect();
                    if (_mrect.width > 0) {
                      const _pp = node.parentElement || node;
                      const _ps = getComputedStyle(_pp);
                      window.__measCtx = window.__measCtx
                        || document.createElement('canvas').getContext('2d');
                      window.__measCtx.font = _ps.fontStyle + ' ' + _ps.fontWeight
                        + ' ' + _ps.fontSize + ' ' + _ps.fontFamily;
                      const _sw = window.__measCtx.measureText('\\u00A0').width
                        || window.__measCtx.measureText(' ').width
                        || (parseFloat(_ps.fontSize) * 0.3);
                      const _n = Math.max(1, Math.round(_mrect.width / Math.max(1, _sw)));
                      runs.push({
                        text: '\\u00A0'.repeat(_n),
                        font_family: _ps.fontFamily,
                        font_size_px: parseFloat(_ps.fontSize),
                        font_weight: _ps.fontWeight,
                        font_style: _ps.fontStyle,
                        color: toRGB(_ps.color),
                        letter_spacing_px: 0,
                        text_transform: 'none',
                        vertical_align: 'baseline',
                        bg_color: '',
                      });
                    }
                  }
                  continue;
                } else {
                  collectRuns(node, runs);
                }
              }
            }
            const after = _pseudoRun(el, '::after');
            if (after) runs.push(after);
          };

          // Walk all elements; classify
          const all = document.querySelectorAll('*');
          for (const el of all) {
            const cs = getComputedStyle(el);
            if (cs.display === 'none' || cs.visibility === 'hidden' ||
                parseFloat(cs.opacity) < 0.01) continue;
            const r = el.getBoundingClientRect();
            if (r.width < 1 || r.height < 1) continue;

            const tag = el.tagName.toLowerCase();

            // Walk depth + detect if any ancestor is a heading (h1-h6).
            // Inside headings, span/etc may be decorative markers (numbered
            // circles, color chips, etc.) that PPT can't replicate as
            // shapes without overlapping the heading text — caller decides
            // to skip those decorations.
            let depth = 0;
            let n = el;
            let inside_heading = false;
            while (n.parentElement) {
              depth++;
              const pt = n.parentElement.tagName.toLowerCase();
              if (pt === 'h1' || pt === 'h2' || pt === 'h3'
                  || pt === 'h4' || pt === 'h5' || pt === 'h6') {
                inside_heading = true;
              }
              n = n.parentElement;
            }

            // Add to elements list (for boxes / images)
            elements.push({
              tag, depth,
              x: r.left, y: r.top, w: r.width, h: r.height,
              src: tag === 'img' ? el.src : null,
              alt: tag === 'img' ? el.alt : null,
              nat_w: tag === 'img' ? (el.naturalWidth || r.width) : null,
              nat_h: tag === 'img' ? (el.naturalHeight || r.height) : null,
              object_fit: tag === 'img' ? cs.objectFit : null,
              bg_color: toRGB(cs.backgroundColor),
              bg_image: cs.backgroundImage,
              // url() background images (QR tiles, icon chips) are painted by
              // the walker as a Picture layered over the element's tile. Capture
              // the sizing keywords so the picture placement matches the browser.
              bg_size: cs.backgroundSize,
              bg_position: cs.backgroundPosition,
              bg_origin: cs.backgroundOrigin,
              border_color: toRGB(cs.borderTopColor),
              border_width_px: parseFloat(cs.borderTopWidth) || 0,
              // Per-side borders — for elements with asymmetric borders
              // (e.g. h2 with only border-bottom, venue-badge with only
              // border-right). Each side rendered as a separate thin Rect
              // in PPT. The legacy border_color/border_width above stays
              // for backwards compat (still represents the top edge).
              border_top_w: parseFloat(cs.borderTopWidth) || 0,
              border_right_w: parseFloat(cs.borderRightWidth) || 0,
              border_bottom_w: parseFloat(cs.borderBottomWidth) || 0,
              border_left_w: parseFloat(cs.borderLeftWidth) || 0,
              border_top_c: toRGB(cs.borderTopColor),
              border_right_c: toRGB(cs.borderRightColor),
              border_bottom_c: toRGB(cs.borderBottomColor),
              border_left_c: toRGB(cs.borderLeftColor),
              border_radius_px: parseFloat(cs.borderTopLeftRadius) || 0,
              box_shadow: cs.boxShadow,
              // Padding: needed for <img> backed by a decorated tile
              // (e.g. `.logo { background:white; padding:10pt 14pt }`).
              // Picture must render INSIDE the padded content area, not
              // inside the bbox edge, or it visually breaks out of the tile.
              pad_top_px: parseFloat(cs.paddingTop) || 0,
              pad_right_px: parseFloat(cs.paddingRight) || 0,
              pad_bottom_px: parseFloat(cs.paddingBottom) || 0,
              pad_left_px: parseFloat(cs.paddingLeft) || 0,
              tag_class: typeof el.className === 'string' ? el.className : '',
              inside_heading,
              is_inline_display: (cs.display === 'inline' || cs.display === 'inline-block'),
              // Parent display lets the Python side detect chip-pill patterns
              // (inline element with bg + significant border-radius inside a
              // flex/grid container — e.g. .p-chips span). When detected, the
              // chip is rendered as its own RoundedRect Shape with text in
              // the shape's text_frame, rather than as a run-level highlight
              // (which can't preserve the inter-pill gap or the rounded
              // corners — the visible bug a user reported on
              // structured-pruning-survey poster.pptx).
              parent_display: el.parentElement ? getComputedStyle(el.parentElement).display : '',
              // Text styling captured per element for the chip-pill rendering
              // path. Python needs these to populate the shape's text_frame
              // since chip spans don't qualify for the normal text_blocks
              // emit (a <span> isn't in BLOCK_TAGS and isn't 'div').
              text_color: toRGB(cs.color),
              font_family: cs.fontFamily,
              font_size_px: parseFloat(cs.fontSize),
              // chip-pill text needs the element's text-transform so a pill
              // like `.tag {text-transform:uppercase}` renders "DECISIVE" not
              // the raw DOM "Decisive" (the chip path bypasses the run pipeline).
              text_transform: cs.textTransform,
              font_weight: cs.fontWeight,
              font_style: cs.fontStyle,
              inner_text: el.textContent.trim(),
              // Element-child count (children = ELEMENT nodes only, ignoring
              // text nodes). Lets the Python chip-pill detector require
              // "leaf" semantics — only TRUE chip pills (e.g. <span class=
              // "tag">Blocked</span>) qualify, NOT composite tiles like
              // <div class="cell"><div class="v">16 GB</div><div class="l">
              // GPT-2 corpus</div></div>. Without this check the chip
              // detector matched the .cell (radius+flex-parent+short total
              // text) and injected the concatenated children text ("16 GB
              // GPT-2 corpus") into the rect, AND the children .v/.l also
              // emitted their own textboxes → 2× duplicate text per cell.
              child_element_count: el.children.length,
            });

            // Pseudo-element decoration emission (issue: P8 numbered-step
            // circles disappeared in .pptx). For ::before / ::after that
            // have visible chrome (bg + non-trivial border-radius + non-zero
            // width/height), emit a SYNTHETIC chip-pill-shaped element so
            // the Python chip-pill detector renders it as a RoundedRect
            // shape with the resolved content (e.g. "1", "2", "3" for a
            // CSS counter) injected into the shape's text_frame. Without
            // this the pseudo's bg/round-shape was dropped entirely (the
            // text run was emitted but inherited no decoration).
            //
            // P8 step-badge fix: when a leading inline ::before marker badge
            // (e.g. .p-steps .step counter) is emitted as its own shape, record
            // how far the following text must shift right to clear it.
            let _beforeMarkerOffsetPx = null;
            const _emitPseudoDeco = (which) => {
              const pcs = getComputedStyle(el, which);
              if (!pcs) return;
              if (pcs.display === 'none' || pcs.visibility === 'hidden') return;
              const raw = pcs.content;
              if (!raw || raw === 'none' || raw === 'normal') return;
              const bg = toRGB(pcs.backgroundColor);
              if (!bg || bg === 'transparent' || bg.includes(', 0)')) return;
              const wPx = parseFloat(pcs.width) || 0;
              const hPx = parseFloat(pcs.height) || 0;
              if (wPx < 4 || hPx < 4) return;
              const radiusStr = pcs.borderTopLeftRadius || '0';
              const radiusVal = parseFloat(radiusStr) || 0;
              // Resolve CSS content: strip quotes, decode \\uXXXX, resolve
              // counter() (Chromium returns the literal 'counter(name)'
              // from getComputedStyle, NOT the resolved number — we must
              // walk the DOM to compute counter-reset / counter-increment).
              let text = raw;
              if ((text.startsWith('"') && text.endsWith('"')) ||
                  (text.startsWith("'") && text.endsWith("'"))) {
                text = text.slice(1, -1);
              }
              text = text.replace(/\\\\([0-9a-fA-F]{1,6})\\s?/g,
                (_, hex) => String.fromCodePoint(parseInt(hex, 16)));
              text = text.replace(
                /counter\\(\\s*([a-zA-Z_][\\w-]*)\\s*(?:,\\s*[^)]+)?\\s*\\)/g,
                (_, name) => {
                  let cnt = 0;
                  let found = false;
                  const walkCnt = (n) => {
                    if (found) return;
                    const ncs = getComputedStyle(n);
                    if (ncs.counterReset && ncs.counterReset !== 'none'
                        && ncs.counterReset.includes(name)) cnt = 0;
                    if (ncs.counterIncrement && ncs.counterIncrement !== 'none'
                        && ncs.counterIncrement.includes(name)) cnt += 1;
                    // Pseudo-elements often carry the counter-increment
                    // (e.g. `.p-steps .step::before { counter-increment: step }`)
                    // — check both pseudos. Per CSS spec, ::before's increment
                    // applies BEFORE its content is laid out, so the value
                    // available to counter() inside content is the post-
                    // increment value.
                    for (const pseudo of ['::before', '::after']) {
                      const pcs2 = getComputedStyle(n, pseudo);
                      if (pcs2 && pcs2.counterIncrement
                          && pcs2.counterIncrement !== 'none'
                          && pcs2.counterIncrement.includes(name)) cnt += 1;
                    }
                    if (n === el) { found = true; return; }
                    for (const ch of n.children) {
                      walkCnt(ch);
                      if (found) return;
                    }
                  };
                  walkCnt(document.body);
                  return String(cnt);
                });
              if (!text.trim()) return;
              const padLeft = parseFloat(cs.paddingLeft) || 0;
              const padTop  = parseFloat(cs.paddingTop) || 0;
              let pseudoX = r.left + padLeft;
              let pseudoY = r.top + padTop;
              if (cs.display === 'flex' || cs.display === 'inline-flex') {
                if (cs.alignItems === 'center') {
                  pseudoY = r.top + Math.max(0, (r.height - hPx) / 2);
                }
                if (cs.justifyContent === 'center') {
                  pseudoX = r.left + Math.max(0, (r.width - wPx) / 2);
                }
              }
              // Normalize % border-radius to px (for `border-radius: 50%`
              // on a circle, half-the-min-dim gives the right adj for PPT).
              const radiusPx = radiusStr.endsWith('%')
                ? Math.round(radiusVal / 100 * Math.min(wPx, hPx))
                : radiusVal;
              if (which === '::before'
                  && cs.display !== 'flex' && cs.display !== 'inline-flex'
                  && (pcs.position === 'static' || pcs.position === 'relative')
                  && (pcs.display || '').includes('inline')) {
                // Block-flow inline ::before (P8 .p-steps .step badge) pushes
                // the following text right by its width + margin-right; mirror
                // that for the separate text box so they don't overlap.
                _beforeMarkerOffsetPx = padLeft + wPx + (parseFloat(pcs.marginRight) || 0);
              }
              elements.push({
                tag: 'span', depth: depth + 1,
                x: pseudoX, y: pseudoY, w: wPx, h: hPx,
                src: null, alt: null, nat_w: null, nat_h: null, object_fit: null,
                bg_color: bg,
                bg_image: pcs.backgroundImage,
                border_color: toRGB(pcs.borderTopColor),
                border_width_px: parseFloat(pcs.borderTopWidth) || 0,
                border_top_w: parseFloat(pcs.borderTopWidth) || 0,
                border_right_w: parseFloat(pcs.borderRightWidth) || 0,
                border_bottom_w: parseFloat(pcs.borderBottomWidth) || 0,
                border_left_w: parseFloat(pcs.borderLeftWidth) || 0,
                border_top_c: toRGB(pcs.borderTopColor),
                border_right_c: toRGB(pcs.borderRightColor),
                border_bottom_c: toRGB(pcs.borderBottomColor),
                border_left_c: toRGB(pcs.borderLeftColor),
                border_radius_px: radiusPx,
                box_shadow: '',
                pad_top_px: parseFloat(pcs.paddingTop) || 0,
                pad_right_px: parseFloat(pcs.paddingRight) || 0,
                pad_bottom_px: parseFloat(pcs.paddingBottom) || 0,
                pad_left_px: parseFloat(pcs.paddingLeft) || 0,
                tag_class: '_pseudo_decoration',
                inside_heading,
                is_inline_display: false,
                // 'flex' makes the chip-pill detector fire (line ~1276 expects
                // parent_display in flex/grid). Even though there's no real
                // parent element here, the pseudo always conceptually sits
                // within its host's flex/grid context.
                parent_display: 'flex',
                text_color: toRGB(pcs.color),
                font_family: pcs.fontFamily,
                font_size_px: parseFloat(pcs.fontSize),
                font_weight: pcs.fontWeight,
                font_style: pcs.fontStyle,
                inner_text: text,
              });
            };
            _emitPseudoDeco('::before');
            _emitPseudoDeco('::after');

            // Text block: this is a BLOCK_TAG with inline-only descendants,
            // OR a leaf element with text + no element children.
            // Skip nested block tags — only the OUTER block emits the text;
            // descendants would re-emit and overlap.
            const has_text = el.textContent.trim().length > 0;
            // Chip-container detector. A flex/grid container whose children
            // will be rendered individually as RoundedRect chip-pills (in
            // the Python build_pptx decoration-box pass). isInlineOnly
            // returns true for these containers because span/em/strong are
            // semantically inline-flow regardless of computed display, so
            // without this check the .p-chips div would ALSO emit a text
            // block with all chip labels concatenated — drawing the chip
            // text twice (once inside each pill shape, once as overlapping
            // textbox). Same suppression for inline-flex / grid /
            // inline-grid parents that use the chip-list pattern.
            const isChipContainer = (() => {
              if (!['flex', 'inline-flex', 'grid', 'inline-grid']
                    .includes(cs.display)) return false;
              let chipCount = 0, nonChipWithTextCount = 0;
              for (const c of el.children) {
                const ccs = getComputedStyle(c);
                // A display:none / visibility:hidden child paints nothing and
                // must NOT make its parent a "chip container". (Legend-frame
                // .section h2 is display:flex and contains a hidden pill-styled
                // .listen-btn "Listen"; without this skip the button counts as
                // the sole chip and the real title text gets suppressed.)
                if (ccs.display === 'none' || ccs.visibility === 'hidden'
                    || parseFloat(ccs.opacity || '1') === 0) continue;
                const radius = parseFloat(ccs.borderTopLeftRadius) || 0;
                const bg = toRGB(ccs.backgroundColor);
                const ctxt = c.textContent.trim();
                // Thresholds MUST mirror the Python chip-pill detector below.
                const isChip = (radius >= 4 && bg && bg !== 'transparent'
                    && !bg.includes(', 0)') && c.children.length === 0
                    && ctxt.length > 0 && ctxt.length <= 60);
                if (isChip) {
                  chipCount += 1;
                } else if (ctxt.length > 0) {
                  nonChipWithTextCount += 1;
                }
              }
              // Only suppress parent text_block when EVERY text-bearing
              // child is a chip. P10 .p-chips (all spans are chips) →
              // suppress ✓. P17 .p-banner (tag chip + body span text) →
              // DON'T suppress, or body text dies. The collectRuns walker
              // separately skips chip-pill descendants, so parent emits
              // ONLY the non-chip text in that case.
              return chipCount > 0 && nonChipWithTextCount === 0;
            })();
            // Chip-pill SELF check. The Python chip-pill detector (build_pptx)
            // renders an element with bg + radius + short text + flex/grid
            // parent as a RoundedRect Shape with the text injected into its
            // text_frame. If we ALSO emit this element as a text_block, the
            // Pass-2 textbox draws the same text on top of the chip's own
            // centered text — visible as overlapping/bold-like double text
            // (the original bug for .p-timeline-stepcards .step .n: a 1.6em
            // circle div with "1"/"2"/"3"). Suppress the text_block emission
            // when this element WILL be picked up as a chip-pill.
            // Thresholds MUST mirror the Python detector (line ~1272).
            const isChipPillItself = (() => {
              const radius = parseFloat(cs.borderTopLeftRadius) || 0;
              if (radius < 4) return false;
              if (el.children.length > 0) return false;
              const bg = toRGB(cs.backgroundColor);
              if (!bg || bg === 'transparent' || bg.includes(', 0)')) return false;
              const txt = el.textContent.trim();
              if (txt.length === 0 || txt.length > 60) return false;
              const parent_disp = el.parentElement
                ? getComputedStyle(el.parentElement).display : '';
              if (!['flex', 'inline-flex', 'grid', 'inline-grid']
                    .includes(parent_disp)) return false;
              if (!['span', 'div'].includes(tag)) return false;
              if (r.height > 150 || r.width > 600) return false;
              return true;
            })();
            // A standalone <span> label that flexbox/grid blockified into its
            // own box (e.g. the scan-to-read QR tile's `<span>Project</span>`
            // caption, which sits under a sibling `.qr-img` div). It is not in
            // BLOCK_TAGS so should_emit_text would drop it, and its parent .tq
            // is NOT inline-only (it holds a block .qr-img child) so the parent
            // never emits the label either — the caption vanished from the pptx.
            // Emit it as its own text block. Guarded so we never DOUBLE-emit:
            // require a non-inline-only parent (else the parent already carries
            // this span's text) and never pull number/kicker spans out of an
            // h1-h6 heading (headings emit their own combined text block).
            const isFlexLabelSpan = (() => {
              if (tag !== 'span') return false;
              if (el.children.length > 0) return false;   // leaf text only
              if (!has_text) return false;
              if (isChipPillItself) return false;         // pills → shape+text
              const par = el.parentElement;
              if (!par) return false;
              const pd = getComputedStyle(par).display;
              if (!['flex', 'inline-flex', 'grid', 'inline-grid'].includes(pd)) return false;
              // Never pull number/kicker spans out of an h1-h6 heading (headings
              // emit their own combined text block).
              for (let p = par; p; p = p.parentElement) {
                if (['h1','h2','h3','h4','h5','h6'].includes(p.tagName.toLowerCase())) return false;
              }
              // Walk up to decide who owns this span's text. If we reach a
              // composite row (e.g. the header .strip — which is now skipped for
              // text emission) before any inline-only block/div ancestor, the row
              // will NOT emit our text, so this leaf must emit itself (the venue
              // fallback label + the QR caption both land here). If instead an
              // inline-only block/div ancestor comes first, IT already carries
              // this text — don't duplicate.
              for (let p = par; p && p !== document.body; p = p.parentElement) {
                if (isCompositeRow(p)) return true;
                const t = p.tagName.toLowerCase();
                if ((BLOCK_TAGS.has(t) || t === 'div') && isInlineOnly(p)) return false;
              }
              return true;
            })();
            const should_emit_text =
              has_text &&
              (BLOCK_TAGS.has(tag) || tag === 'div' || isFlexLabelSpan) &&
              isInlineOnly(el) &&
              !isChipContainer &&
              !isChipPillItself &&
              !isCompositeRow(el);

            if (should_emit_text) {
              const runs = [];
              collectRuns(el, runs);
              // Filter empty runs
              const non_empty_runs = runs.filter(r => r.text.length > 0);
              if (non_empty_runs.length === 0) continue;
              // Browser-mimicking whitespace collapse. HTML treats
              // whitespace-only text nodes between element children as
              // collapsible (a run of newline+indent renders as one
              // space, or nothing at all at block boundaries). PPT text
              // frames keep every newline as a paragraph break and every
              // space as literal text. Without trimming here we get
              // leading blank paragraphs from the indentation around a
              // skipped mjx-container or chip-pill — concrete repro on
              // ResNet .p-eq left a ~3 inch gap above the body text.
              const _isWS = (r) => !/[^ \\t\\r\\n]/.test(r.text);
              let _start = 0, _end = non_empty_runs.length;
              while (_start < _end && _isWS(non_empty_runs[_start])) _start += 1;
              while (_end > _start && _isWS(non_empty_runs[_end - 1])) _end -= 1;
              const _trimmed = non_empty_runs.slice(_start, _end);
              const collapsed_runs = [];
              for (const r of _trimmed) {
                if (_isWS(r) && collapsed_runs.length > 0
                    && _isWS(collapsed_runs[collapsed_runs.length - 1])) {
                  continue;
                }
                if (_isWS(r)) {
                  collapsed_runs.push({...r, text: ' '});
                } else {
                  collapsed_runs.push(r);
                }
              }
              if (collapsed_runs.length === 0) continue;

              // For <li>: pull parent ul/ol's padding-left to know how far
              // left to extend the textbox for the bullet (so bullet sits
              // in the list's marker area, not at li content edge).
              let parent_list_pad_px = 0;
              if (tag === 'li') {
                let p = el.parentElement;
                while (p && !['ul','ol'].includes(p.tagName.toLowerCase())) {
                  p = p.parentElement;
                }
                if (p) {
                  parent_list_pad_px = parseFloat(getComputedStyle(p).paddingLeft) || 0;
                }
              }

              // Tight-bbox recompute. If this element has descendants
              // that the walker skipped (mjx-container math, or chip-
              // pill spans), the parent's getBoundingClientRect covers
              // the WHOLE area including those skipped child regions —
              // but the surviving text in collapsed_runs only occupies
              // a small slice (typically a description below the math,
              // or text around chip pills). Using the parent's bbox
              // makes the textbox huge with the text anchored at the
              // top, drawing the description over the equation area.
              // Recompute the tight bbox via Range API over surviving
              // text nodes, and pull line-height from the closest
              // ancestor of those nodes (so .where's 1.0× spacing is
              // honored instead of .p-eq's larger inherited value).
              let tb_x = r.left, tb_y = r.top, tb_w = r.width, tb_h = r.height;
              let tb_lh = cs.lineHeight;
              const _skippable = (n) => {
                if (n.nodeType !== Node.ELEMENT_NODE) return false;
                const tn = n.tagName.toLowerCase();
                return tn === 'mjx-container'
                    || (n.classList && n.classList.contains('katex'))
                    || isChipPillNode(n);
              };
              let _hasSkipped = false;
              for (const cand of el.querySelectorAll('*')) {
                if (_skippable(cand)) { _hasSkipped = true; break; }
              }
              if (_hasSkipped) {
                const _survivors = [];
                const _walk = (n) => {
                  if (n.nodeType === Node.ELEMENT_NODE) {
                    if (_skippable(n)) return;
                    for (const c of n.childNodes) _walk(c);
                  } else if (n.nodeType === Node.TEXT_NODE
                             && n.textContent.replace(/[ \\t\\r\\n]/g, '').length > 0) {
                    _survivors.push(n);
                  }
                };
                _walk(el);
                if (_survivors.length > 0) {
                  const _rg = document.createRange();
                  let mL = Infinity, mT = Infinity, mR = -Infinity, mB = -Infinity;
                  for (const n of _survivors) {
                    _rg.selectNodeContents(n);
                    const _rc = _rg.getBoundingClientRect();
                    if (_rc.width > 0 && _rc.height > 0) {
                      if (_rc.left   < mL) mL = _rc.left;
                      if (_rc.top    < mT) mT = _rc.top;
                      if (_rc.right  > mR) mR = _rc.right;
                      if (_rc.bottom > mB) mB = _rc.bottom;
                    }
                  }
                  if (isFinite(mL) && isFinite(mT)) {
                    tb_x = mL; tb_y = mT; tb_w = mR - mL; tb_h = mB - mT;
                    const _inner = _survivors[0].parentElement;
                    if (_inner) tb_lh = getComputedStyle(_inner).lineHeight;
                  }
                }
              }

              if (_beforeMarkerOffsetPx != null && tb_x <= r.left + 1) {
                // The leading ::before marker badge was emitted as its own
                // shape at the element's left edge (a PPT text frame can't hold
                // it — runs have no solid background). Shift the text box past
                // the badge + its margin-right gap so the number and the text
                // don't pile up at the same start position.
                const _padR = parseFloat(cs.paddingRight) || 0;
                tb_x = r.left + _beforeMarkerOffsetPx;
                tb_w = Math.max(10, r.width - _beforeMarkerOffsetPx - _padR);
              }

              // Flex/grid COLUMN container whose inline children stack
              // vertically as flex items (e.g. .scan-link, a flex-column div
              // holding span.lk-label "Project" over span.lk-url URL). It gets
              // force-wrapped into multiple paragraphs (see the Python gate).
              const _isFlexCol = (cs.display === 'flex' || cs.display === 'inline-flex')
                                 && (cs.flexDirection || '').startsWith('column');
              // The CONTAINER's inherited line-height is often huge (a large
              // section font-size cascades down); applied to the force-wrapped
              // paragraphs it spaces them far apart (observed: the scan-link
              // "Project" / URL pair split by ~37pt of dead space). Use the
              // first flex child's OWN line-height so the pair reads tight.
              let _lineH = tb_lh;
              if (_isFlexCol && el.firstElementChild) {
                const _clh = getComputedStyle(el.firstElementChild).lineHeight;
                if (_clh && _clh !== 'normal') _lineH = _clh;
              }

              text_blocks.push({
                tag, depth,
                x: tb_x, y: tb_y, w: tb_w, h: tb_h,
                text_align: cs.textAlign,
                line_height: _lineH,
                text_transform: cs.textTransform,
                pad_top_px: parseFloat(cs.paddingTop) || 0,
                pad_right_px: parseFloat(cs.paddingRight) || 0,
                pad_bottom_px: parseFloat(cs.paddingBottom) || 0,
                pad_left_px: parseFloat(cs.paddingLeft) || 0,
                overflow_hidden: cs.overflow === 'hidden' || cs.overflowY === 'hidden',
                webkit_line_clamp: parseInt(cs.webkitLineClamp || '0') || 0,
                is_list_item: tag === 'li',
                parent_list_pad_px,
                hyphens: cs.hyphens || cs.webkitHyphens || 'manual',
                wrap_lines: extractWrapLines(el),
                // PPT has no flex layout, so a flex-column container's stacked
                // spans would render on ONE line ("Projectaka.ms/...") without
                // explicit newlines. This flag makes the Python force-wrap gate
                // honor the browser per-item Y splits (extractWrapLines already
                // measured them) so the spans stack as separate lines.
                is_flex_col: _isFlexCol,
                runs: collapsed_runs,
              });
            }
          }

          // === Orphan text-node scan (Bug 2 backstop) ===
          // The main loop above emits text only from elements where ALL
          // children are inline (isInlineOnly check). When a block element
          // mixes block children with raw text nodes — e.g.
          //   <div class="side"><h4>Title</h4>Body text.</div>
          // — "Body text." is a raw text node directly under .side. The
          // parent .side has a <h4> child (block), so isInlineOnly returns
          // false and .side isn't emitted; the orphan text node has no
          // element of its own, so it's lost.
          //
          // Recipe-level fix is to wrap such text in <p> (the .p-vs recipe
          // in content_patterns.md now mandates this). This walker scan is
          // the defensive backstop: walk every element with a text-node
          // child sitting alongside block-element siblings, emit each such
          // text node as its own text block via Range API bbox.
          //
          // ONLY runs on prose-block containers (div / section / article /
          // aside / blockquote). Headings (h1-h6), interactive elements
          // (button / a / label), and list constructs (ul / ol / li / dl)
          // have their own emission paths — e.g. an h2 with `display:flex`
          // and a <button class="listen-btn"> child has the heading text
          // "Method" as an orphan text node next to the block-displayed
          // button (button is flexified to display:block). The h2 already
          // emits "Method Listen" as one text block via the normal path;
          // re-emitting "Method" alone via the orphan scan would draw the
          // word twice on top of itself.
          const ORPHAN_PROSE_TAGS = new Set([
            'div', 'section', 'article', 'aside', 'blockquote', 'main', 'figure',
          ]);
          for (const el of all) {
            if (!ORPHAN_PROSE_TAGS.has(el.tagName.toLowerCase())) continue;
            const cs = getComputedStyle(el);
            if (cs.display === 'none' || cs.visibility === 'hidden') continue;
            // Only consider mixed-block-and-text parents.
            let hasBlockChild = false;
            let hasOrphanText = false;
            for (const child of el.childNodes) {
              if (child.nodeType === Node.ELEMENT_NODE) {
                const ccs = getComputedStyle(child);
                if (!(ccs.display === 'inline' || ccs.display === 'inline-block')) {
                  hasBlockChild = true;
                }
              } else if (child.nodeType === Node.TEXT_NODE
                         && child.textContent.trim().length > 0) {
                hasOrphanText = true;
              }
            }
            if (!(hasBlockChild && hasOrphanText)) continue;

            for (const child of el.childNodes) {
              if (child.nodeType !== Node.TEXT_NODE) continue;
              const text = child.textContent.trim();
              if (text.length === 0) continue;
              const range = document.createRange();
              range.selectNode(child);
              const r = range.getBoundingClientRect();
              if (r.width < 1 || r.height < 1) continue;
              text_blocks.push({
                tag: 'orphan-text',
                depth: 0,
                x: r.left, y: r.top, w: r.width, h: r.height,
                text_align: cs.textAlign,
                line_height: cs.lineHeight,
                text_transform: cs.textTransform,
                pad_top_px: 0, pad_right_px: 0,
                pad_bottom_px: 0, pad_left_px: 0,
                overflow_hidden: false,
                webkit_line_clamp: 0,
                is_list_item: false,
                parent_list_pad_px: 0,
                hyphens: cs.hyphens || cs.webkitHyphens || 'manual',
                wrap_lines: [],
                runs: [{
                  text: text,
                  font_family: cs.fontFamily,
                  font_size_px: parseFloat(cs.fontSize),
                  font_weight: cs.fontWeight,
                  font_style: cs.fontStyle,
                  color: toRGB(cs.color),
                  letter_spacing_px: cs.letterSpacing === 'normal'
                                     ? 0
                                     : (parseFloat(cs.letterSpacing) || 0),
                  text_transform: cs.textTransform,
                  vertical_align: 'baseline',
                  bg_color: '',
                }],
              });
            }
          }

          // === Math pass: capture rendered equations as native PowerPoint
          // OMML candidates. Engine-agnostic — supports BOTH renderers so the
          // template can migrate MathJax<->KaTeX without touching this path:
          //   * MathJax: <mjx-container> stamped with data-tex (the original
          //     TeX) + data-display ('1' for display, '0' for inline) by the
          //     template's startup.ready hook.
          //   * KaTeX: <span class="katex"> with the source TeX in a nested
          //     <annotation encoding="application/x-tex">, display math wrapped
          //     in <span class="katex-display">.
          // Python Pass 3 turns each `tex` into <a14:m><m:oMath> (editable math
          // in PowerPoint) and prefers the rasterized PNG (below) for soffice.
          // Older templates without either stamp leave `tex` empty and fall
          // back to the SVG-as-image path silently.
          const math_blocks = [];
          const _pushMath = (el, tex, display) => {
            if (!tex) return;
            const r = el.getBoundingClientRect();
            if (r.width === 0 || r.height === 0) return;
            const cs = getComputedStyle(el);
            math_blocks.push({
              x: r.left, y: r.top, w: r.width, h: r.height,
              tex,
              display,
              color: toRGB(cs.color),
              font_size_px: parseFloat(cs.fontSize),
            });
          };
          // MathJax
          for (const c of document.querySelectorAll('mjx-container')) {
            _pushMath(c, c.getAttribute('data-tex'),
                      c.getAttribute('data-display') === '1');
          }
          // KaTeX (read TeX from the a11y MathML annotation; skip the hidden
          // .katex-mathml mirror so each equation is captured exactly once)
          for (const k of document.querySelectorAll('.katex')) {
            if (k.closest('mjx-container')) continue;
            const _ann = k.querySelector('annotation[encoding="application/x-tex"]');
            const _tex = (_ann && _ann.textContent)
                         || k.getAttribute('data-tex') || '';
            _pushMath(k, _tex.trim(), k.closest('.katex-display') !== null);
          }

          return {viewport_w: window.innerWidth, viewport_h: window.innerHeight,
                  body_w: document.body.scrollWidth,
                  body_h: document.body.scrollHeight,
                  elements, text_blocks, math_blocks};
        }""")
        # Rasterize each equation region as a PNG fallback while the page is open.
        # PowerPoint renders the native OMML we emit in Pass 3, but LibreOffice/
        # soffice (and some other viewers) do NOT rasterize the <a14:m> slide-math
        # extension, leaving the equation BLANK. A screenshot of the browser-
        # rendered MathJax keeps it visible everywhere; Pass 3 prefers it over OMML.
        try:
            import tempfile as _tf, os as _os
            _mroot = _tf.mkdtemp(prefix="p2p_math_")
            for _i, _mb in enumerate(data.get("math_blocks") or []):
                try:
                    if _mb.get("w", 0) < 1 or _mb.get("h", 0) < 1:
                        continue
                    _pp = _os.path.join(_mroot, f"eq{_i}.png")
                    page.screenshot(path=_pp, clip={
                        "x": max(0.0, float(_mb["x"])), "y": max(0.0, float(_mb["y"])),
                        "width": max(1.0, float(_mb["w"])), "height": max(1.0, float(_mb["h"]))})
                    _mb["png"] = _pp
                except Exception:
                    pass
        except Exception:
            pass
        b.close()
    return data


# ─── shape generation ──────────────────────────────────────────────────────

def has_decoration(el: dict) -> bool:
    """Element has visual styling (bg / border / gradient) worth emitting as a Rectangle."""
    if _parse_color(el["bg_color"]):
        return True
    # Check all 4 borders (asymmetric border support)
    for side in ("top", "right", "bottom", "left"):
        w = el.get(f"border_{side}_w", 0)
        c = el.get(f"border_{side}_c") if f"border_{side}_c" in el else el.get("border_color")
        if w > 0 and _parse_color(c):
            return True
    # Legacy fallback
    if el["border_width_px"] > 0 and _parse_color(el["border_color"]):
        return True
    if el["bg_image"] and "gradient" in (el["bg_image"] or ""):
        return True
    return False


def is_full_page(el: dict, body_w: float, body_h: float) -> bool:
    """Skip global page-bg elements (html, body, .poster wrapper)."""
    return el["tag"] in ("html", "body") or (
        el["x"] <= 1 and el["y"] <= 1 and
        el["w"] >= body_w * 0.99 and el["h"] >= body_h * 0.99)


def _resolve_bg_url(bg_image: str) -> str | None:
    """Extract a readable local file path from a CSS `background-image:url(...)`.

    The scan-to-read QR tiles (and some CSS icon chips) paint their image via a
    div `background-image` rather than an <img>, so the figure-sizing / preflight
    gates never treat them as figures. The walker must still rasterize them, or
    they vanish in the pptx (the visible "empty white QR box" symptom).

    Returns None for gradients / data-URIs / unreadable paths. Chromium resolves
    a relative url() to an absolute `file://` when the page is loaded over
    file://, so we unwrap the scheme and percent-decode."""
    if not bg_image or "url(" not in bg_image or "gradient" in bg_image:
        return None
    m = re.search(r'url\((["\']?)(.*?)\1\)', bg_image)
    if not m:
        return None
    u = m.group(2).strip()
    if u.startswith("data:"):
        return None  # data-URI backgrounds are rare here; not handled
    if u.startswith("file://"):
        u = u[7:]
    from urllib.parse import unquote
    u = unquote(u)
    return u if u and Path(u).exists() else None


def build_pptx(dom: dict, out_path: Path,
               width_inch: float = DEFAULT_W_INCH,
               height_inch: float = DEFAULT_H_INCH,
               corrections: dict | None = None,
               match_wrap: bool = False) -> dict:
    """corrections: {block_idx_str: {"font_scale": float}} — per-text-block
    font-size scale factor applied by L2 closed-loop. Index = position in
    sorted text_blocks list.

    slide_scale: when the design canvas (e.g. 60×36") exceeds PPT's 56-inch
    per-axis cap and we cap to a smaller slide (e.g. 55×33"), absolute units
    (font size in Pt, line spacing in Pt, line width in EMU, letter-spacing
    in OOXML spc) must scale to match — otherwise text appears proportionally
    larger than designed. Coordinate-space (positions, widths, heights) auto-
    scales via emu_x/emu_y which use the slide_w_emu / body_w_px ratio."""
    corrections = corrections or {}
    prs = Presentation()
    prs.slide_width = Inches(width_inch)
    prs.slide_height = Inches(height_inch)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height
    body_w_px = dom.get("body_w") or dom["viewport_w"]
    body_h_px = dom.get("body_h") or dom["viewport_h"]

    # Detect slide-vs-design scale: design canvas was viewport_w/96 inches.
    # If slide is smaller (cap), scale all absolute units uniformly.
    design_w_inch = (dom.get("viewport_w") or body_w_px) / 96.0
    slide_scale = width_inch / design_w_inch if design_w_inch > 0 else 1.0
    if abs(slide_scale - 1.0) > 0.001:
        print(f"      [scale] design {design_w_inch:.1f}\" → slide {width_inch:.1f}\" "
              f"= {slide_scale:.3f}× (font/line/border auto-scaled)", file=sys.stderr)

    def emu_x(px): return _px_to_emu(px, slide_w_emu, body_w_px)
    def emu_y(px): return _px_to_emu(px, slide_h_emu, body_h_px)

    counts = {"text": 0, "pictures": 0, "boxes": 0, "gradients": 0,
              "shadows": 0, "skipped_full_page": 0}

    # ─── Pass 1: emit IMAGES + DECORATIVE BOXES (z-order: containers below leaves) ──
    elements = sorted(dom["elements"], key=lambda e: (e["depth"], e["y"], e["x"]))
    for el in elements:
        if is_full_page(el, body_w_px, body_h_px):
            counts["skipped_full_page"] += 1
            continue

        x_emu = emu_x(el["x"])
        y_emu = emu_y(el["y"])
        w_emu = emu_x(el["x"] + el["w"]) - x_emu
        h_emu = emu_y(el["y"] + el["h"]) - y_emu

        # Image
        if el["tag"] == "img" and el["src"]:
            src = el["src"]
            # ─ Decoration tile BENEATH the picture ─
            # Logos, badges, icon-boxes often have CSS like
            #   .logo { background: white; padding: 10pt 14pt; border-radius: 8pt; }
            # The picture alone on a colored banner has transparent-PNG
            # text vanish into the background. Emit a backing Rectangle/
            # RoundedRect FIRST (z-order: below picture) so the transparent
            # logo sits on its designed tile, not on the section background.
            if has_decoration(el) and not el.get("is_inline_display"):
                try:
                    radius_px = el.get("border_radius_px", 0)
                    use_round = radius_px >= 2
                    shape_type = (MSO_SHAPE.ROUNDED_RECTANGLE if use_round
                                  else MSO_SHAPE.RECTANGLE)
                    tile = slide.shapes.add_shape(shape_type, x_emu, y_emu, w_emu, h_emu)
                    if use_round:
                        _set_round_rect_adj(tile, radius_px, el["w"], el["h"])
                    bg = _parse_color(el["bg_color"])
                    if bg:
                        tile.fill.solid()
                        tile.fill.fore_color.rgb = bg[0]
                    else:
                        tile.fill.background()
                    bd = _parse_color(el.get("border_color", ""))
                    if bd and el.get("border_width_px", 0) > 0:
                        tile.line.color.rgb = bd[0]
                        tile.line.width = Emu(int(el["border_width_px"] * 9525 * slide_scale))
                    else:
                        tile.line.fill.background()
                    counts["boxes"] += 1
                except Exception as e:
                    print(f"[skip img tile {src[-40:]}] {e}", file=sys.stderr)
            img_path = None
            try:
                if src.startswith("data:"):
                    m = re.match(r"data:image/[^;]+;base64,(.+)", src)
                    if m:
                        with tempfile.NamedTemporaryFile(suffix=".png",
                                                         delete=False) as f:
                            f.write(base64.b64decode(m.group(1)))
                            img_path = f.name
                elif src.startswith("file://"):
                    img_path = src[7:]
                if img_path and Path(img_path).exists():
                    # SVG isnt natively supported by python-pptx (PIL cant
                    # decode it). Rasterize to PNG at 2× density first.
                    if img_path.lower().endswith(".svg"):
                        try:
                            import cairosvg  # type: ignore
                            png_path = tempfile.NamedTemporaryFile(
                                suffix=".png", delete=False).name
                            cairosvg.svg2png(url=img_path, write_to=png_path,
                                              output_width=max(400, int(el.get("nat_w") or 400) * 2))
                            img_path = png_path
                        except ImportError:
                            print(f"[skip svg {src[-40:]}] cairosvg not installed", file=sys.stderr)
                            continue
                        except Exception as e:
                            print(f"[skip svg {src[-40:]}] {e}", file=sys.stderr)
                            continue
                    # object-fit: contain → image is centered inside the
                    # padded CONTENT area (bbox minus padding), preserving
                    # natural aspect ratio. Without padding subtract, a
                    # picture inside a tile would overflow the white tile
                    # edge into the section background.
                    obj_fit = el.get("object_fit") or "fill"
                    nat_w = el.get("nat_w") or el["w"]
                    nat_h = el.get("nat_h") or el["h"]
                    pad_l = el.get("pad_left_px", 0)
                    pad_r = el.get("pad_right_px", 0)
                    pad_t = el.get("pad_top_px", 0)
                    pad_b = el.get("pad_bottom_px", 0)
                    content_x = el["x"] + pad_l
                    content_y = el["y"] + pad_t
                    content_w = max(1, el["w"] - pad_l - pad_r)
                    content_h = max(1, el["h"] - pad_t - pad_b)
                    px_w, px_h, off_x, off_y = content_w, content_h, 0, 0
                    if obj_fit == "contain" and nat_w > 0 and nat_h > 0:
                        nat_ratio = nat_w / nat_h
                        bbox_ratio = content_w / content_h
                        if nat_ratio > bbox_ratio:
                            # wider than content → fit by width, letterbox top/bottom
                            px_w = content_w
                            px_h = content_w / nat_ratio
                            off_y = (content_h - px_h) / 2
                        else:
                            # taller → fit by height, letterbox left/right
                            px_h = content_h
                            px_w = content_h * nat_ratio
                            off_x = (content_w - px_w) / 2
                    picture_x = emu_x(content_x + off_x)
                    picture_y = emu_y(content_y + off_y)
                    picture_w = emu_x(content_x + off_x + px_w) - picture_x
                    picture_h = emu_y(content_y + off_y + px_h) - picture_y
                    slide.shapes.add_picture(img_path, picture_x, picture_y,
                                             width=picture_w, height=picture_h)
                    counts["pictures"] += 1
            except Exception as e:
                print(f"[skip img {el.get('alt') or src[:40]}] {e}", file=sys.stderr)
            finally:
                if src.startswith("data:") and img_path:
                    Path(img_path).unlink(missing_ok=True)
            continue

        # Decorative box
        if has_decoration(el):
            # Skip decorations inside h1-h6 (e.g., span.num circle markers
            # rendered as standalone shapes). They overlap the heading text
            # in PPT since we can't position the marker inside-the-circle
            # like CSS does. The marker text itself is preserved in the
            # heading's text block.
            if el.get("inside_heading"):
                continue
            # Inline elements with bg become RUN-level <a:highlight> properties
            # in Pass 2 (PowerPoint Text Highlight Color), not standalone
            # rectangles. Skip Pass 1 emission to avoid double-painting.
            if el.get("is_inline_display"):
                continue
            # CHIP-PILL detection (Bug 1 fix). The .p-chips widget puts
            # short text labels inside <span>s, but flexbox blockifies the
            # spans so they reach this branch as block-displayed elements.
            # Recognize a chip by: very-rounded border + small bbox + short
            # inner text + flex/grid parent + leaf (no element children).
            # When detected, after rendering the RoundedRect shape we
            # inject the chip's text into its text_frame so the pill shows
            # its label. Without this, the chip renders as an empty bubble
            # and the label gets dropped from a separate Pass-2 textbox
            # (because <span> isn't in BLOCK_TAGS so should_emit_text skips
            # it). The Pass-2 path silently lost all chip text — the
            # user-reported "concatenated text" symptom.
            #
            # Discrimination thresholds (tuned vs section-card false-positives):
            #   border_radius >= 4    → chip pills set 0.25em+ (≈6px+ at 24px
            #                           font); section cards typically square.
            #                           Threshold was 12 originally, but that
            #                           missed small-radius badges like P17
            #                           .p-banner .tag (radius 0.25em ≈ 6px)
            #                           — the "Blocked"/"First" tag rendered
            #                           as an empty red pill in pptx while
            #                           the parent .p-banner emitted both
            #                           tag and body text concatenated on
            #                           top, breaking the banner layout.
            #   h_px <= 150           → chips are ~50-115px tall; section cards 200+
            #   w_px <= 600           → chips span 150-460px (longest: "MobileNet-V1/2/3");
            #                           section cards are full-column width (1000+px)
            #   inner_text len <= 60  → chip labels are short tokens, never paragraphs
            #   child_element_count == 0 → TRUE chip pills are LEAF elements
            #                           (a `<span class="tag">Blocked</span>`),
            #                           NOT composite tiles like
            #                           `<div class="cell"><div class="v">16 GB</div>
            #                           <div class="l">GPT-2 corpus</div></div>`.
            #                           Without this check the .p-stat-strip
            #                           .cell matched (radius 0.5em, parent
            #                           flex, total text <60) and got its
            #                           CONCATENATED children text ("16 GB
            #                           GPT-2 corpus") injected into the
            #                           rect — and the children .v/.l ALSO
            #                           emitted their own textboxes, producing
            #                           visible duplicate text per cell.
            #   tag span/div          → don't try to inject text into <img>/etc
            chip_text = ""
            parent_disp = el.get("parent_display", "")
            inner_text = (el.get("inner_text", "") or "").strip()
            if (el.get("border_radius_px", 0) >= 4
                    and el.get("child_element_count", 0) == 0
                    and el.get("h", 0) <= 150
                    and el.get("w", 0) <= 600
                    and 0 < len(inner_text) <= 60
                    and parent_disp in ("flex", "inline-flex", "grid", "inline-grid")
                    and el.get("tag") in ("span", "div")):
                chip_text = inner_text
                # Honor the pill's CSS text-transform (e.g. .tag {text-transform:
                # uppercase}) — the chip path bypasses the normal run pipeline that
                # applies it, so a source "Decisive" would otherwise render lower
                # case instead of "DECISIVE". (Inlined: _apply_text_transform is
                # defined later in this function, after this Pass-1 loop.)
                _cttf = (el.get("text_transform") or "none")
                if _cttf == "uppercase":
                    chip_text = chip_text.upper()
                elif _cttf == "lowercase":
                    chip_text = chip_text.lower()
                elif _cttf == "capitalize":
                    chip_text = chip_text.title()
            try:
                # Only use ROUNDED_RECTANGLE for visibly-rounded UI;
                # PPT's preset adj defaults to ~16.7% which would balloon
                # CSS 4px → 30px+ on a tall section. We override adj below.
                radius_px = el["border_radius_px"]
                use_round = radius_px >= 2
                shape_type = (MSO_SHAPE.ROUNDED_RECTANGLE if use_round
                              else MSO_SHAPE.RECTANGLE)
                shape = slide.shapes.add_shape(shape_type, x_emu, y_emu, w_emu, h_emu)
                if use_round:
                    _set_round_rect_adj(shape, radius_px, el["w"], el["h"])

                grad_stops = _parse_gradient(el["bg_image"] or "")
                if grad_stops:
                    _set_gradient_fill(shape, grad_stops, angle_deg=90)
                    counts["gradients"] += 1
                else:
                    bg = _parse_color(el["bg_color"])
                    if bg:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = bg[0]
                    else:
                        shape.fill.background()

                bd = _parse_color(el["border_color"])
                # Determine if borders are SYMMETRIC (all four sides same width
                # + same color) — if so, set shape.line uniformly. Otherwise,
                # don't paint shape.line (we'll emit per-side lines below).
                sides = ["top", "right", "bottom", "left"]
                widths = [el.get(f"border_{s}_w", 0) for s in sides]
                colors = [el.get(f"border_{s}_c", el.get("border_color", "")) for s in sides]
                symmetric = (len(set(widths)) == 1 and len(set(colors)) == 1
                              and widths[0] > 0)
                # BASE border: the (width, color) shared by a MAJORITY of sides.
                # A card with a thin hairline box border + ONE thick accent side
                # (e.g. `.section { border: 1px solid #e8e8e2; border-left: 14pt
                # accent }` — the left-bar style / the takeaway callout) is
                # asymmetric, but its 3 thin sides are the card OUTLINE. Drawing
                # each as its own sharp full-width Rect makes the top/bottom
                # hairlines read as spurious horizontal bars flanking the accent
                # side. Instead paint the base border via the shape's own
                # (corner-radius-following) line and emit a per-side Rect ONLY for
                # the odd/thick side(s).
                _pairs = list(zip(widths, colors))
                _base_pair = max(set(_pairs), key=_pairs.count)
                _base_n = _pairs.count(_base_pair)
                _base_w, _base_c = _base_pair
                base_as_line = (not symmetric and _base_n >= 2
                                and 0 < _base_w <= 2.5 and bool(_parse_color(_base_c)))
                if bd and el["border_width_px"] > 0 and symmetric:
                    shape.line.color.rgb = bd[0]
                    shape.line.width = Emu(int(el["border_width_px"] * 9525 * slide_scale))
                elif base_as_line:
                    shape.line.color.rgb = _parse_color(_base_c)[0]
                    shape.line.width = Emu(int(_base_w * 9525 * slide_scale))
                else:
                    shape.line.fill.background()

                # Emit per-side thin Rectangles for ASYMMETRIC borders
                # (h2 with only border-bottom, venue-badge with only
                # border-right, etc.). Each present side becomes a
                # 1px-thin Rectangle of the right color positioned at
                # the corresponding edge. When base_as_line is set, the
                # majority hairline is already painted as the shape's own
                # line, so skip those sides — only overlay the odd/thick one.
                if not symmetric:
                    for side, w_px, c in zip(sides, widths, colors):
                        if w_px <= 0:
                            continue
                        if base_as_line and (w_px, c) == _base_pair:
                            continue
                        bc = _parse_color(c)
                        if not bc:
                            continue
                        # Thin band coordinates in browser px
                        if side == "top":
                            sx, sy, sw, sh = el["x"], el["y"], el["w"], w_px
                        elif side == "bottom":
                            sx, sy, sw, sh = el["x"], el["y"] + el["h"] - w_px, el["w"], w_px
                        elif side == "left":
                            sx, sy, sw, sh = el["x"], el["y"], w_px, el["h"]
                        else:  # right
                            sx, sy, sw, sh = el["x"] + el["w"] - w_px, el["y"], w_px, el["h"]
                        try:
                            ex_x = emu_x(sx)
                            ex_y = emu_y(sy)
                            ex_w = emu_x(sx + sw) - ex_x
                            ex_h = emu_y(sy + sh) - ex_y
                            # Floor the THIN dimension to ~1.3px: a 1px CSS hairline
                            # (table row rules, section dividers) converts to <1px on
                            # a down-scaled slide and soffice/PPT then drop it. Keep
                            # the long dimension as measured.
                            if side in ("top", "bottom"):
                                ex_h = max(ex_h, int(emu_y(sy + 1.3) - ex_y))
                            else:
                                ex_w = max(ex_w, int(emu_x(sx + 1.3) - ex_x))
                            ex_w = max(1, ex_w)
                            ex_h = max(1, ex_h)
                            side_shape = slide.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE, ex_x, ex_y, ex_w, ex_h)
                            side_shape.fill.solid()
                            side_shape.fill.fore_color.rgb = bc[0]
                            side_shape.line.fill.background()
                            # Strip the python-pptx-auto-generated <p:style>
                            # block. It carries <a:lnRef idx="1" schemeClr=
                            # "accent1">, which Mac PowerPoint sometimes
                            # renders as a thin accent-colored outline ON
                            # TOP of <a:noFill/> — making thin dividers
                            # (.headline-hero .supporting top border, the
                            # SO WHAT separator above section conclusions,
                            # etc.) appear visibly darker than the HTML
                            # alpha-composited color. Removing <p:style>
                            # forces PowerPoint to honor only <a:spPr>'s
                            # explicit fill + line, eliminating the
                            # version-dependent overlay.
                            _ps_sp = side_shape._element
                            _ps_style = _ps_sp.find(qn("p:style"))
                            if _ps_style is not None:
                                _ps_sp.remove(_ps_style)
                        except Exception:
                            pass

                bs = el.get("box_shadow", "") or ""
                if bs and bs != "none":
                    try:
                        _add_shadow(shape, blur_px=5, offset_y_px=2, alpha=0.15)
                        counts["shadows"] += 1
                    except Exception:
                        pass

                # Chip-pill text injection. When the decoration we just
                # rendered is an inline chip-pill (caught by the chip_pill
                # exception above), populate the shape's text_frame with the
                # chip label, vertically + horizontally centered, no wrap.
                # Element-level computed styles drive the font (color, size,
                # weight, italic) so the chip's rendered text matches the
                # browser's chip appearance.
                if chip_text:
                    tf = shape.text_frame
                    tf.margin_left = Emu(0)
                    tf.margin_right = Emu(0)
                    tf.margin_top = Emu(0)
                    tf.margin_bottom = Emu(0)
                    tf.word_wrap = False
                    try:
                        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    except Exception:
                        pass
                    p_ = tf.paragraphs[0]
                    p_.alignment = PP_ALIGN.CENTER
                    run = p_.add_run()
                    run.text = chip_text
                    fsz_px = el.get("font_size_px", 14) or 14
                    run.font.size = Pt(fsz_px * 72 / 96 * slide_scale)
                    fam = _pick_font_family(el.get("font_family", "Inter"))
                    if fam:
                        run.font.name = fam
                    fw = el.get("font_weight", "400")
                    try:
                        run.font.bold = int(fw) >= 600
                    except (ValueError, TypeError):
                        run.font.bold = fw in ("bold", "bolder")
                    fst = el.get("font_style", "normal")
                    if fst in ("italic", "oblique"):
                        run.font.italic = True
                    tc = _parse_color(el.get("text_color", ""))
                    if tc:
                        run.font.color.rgb = tc[0]
                    counts.setdefault("chip_pills", 0)
                    counts["chip_pills"] += 1

                counts["boxes"] += 1
            except Exception as e:
                print(f"[skip box {el['tag']}] {e}", file=sys.stderr)

        # ─ url() BACKGROUND IMAGE → Picture (QR tiles, CSS icon chips) ─
        # A div with `background-image:url(...)` (the scan-to-read QR uses this
        # so the figure-sizing gates never see it as an <img>) paints NOTHING in
        # PPT on its own. The decorative-box branch above already drew the tile
        # (white bg + border + radius); now layer the actual image ON TOP, sized
        # per background-size (contain letterboxes, cover center-crops, else
        # fills the content box). <img>/full-page elements never reach here
        # (they `continue` earlier), so this only fires for bg-image divs.
        _bgi = el.get("bg_image") or ""
        if ("url(" in _bgi and "gradient" not in _bgi
                and not el.get("is_inline_display")):
            _bgpath = _resolve_bg_url(_bgi)
            if _bgpath:
                try:
                    from PIL import Image as _PILImage
                    with _PILImage.open(_bgpath) as _im:
                        _nw, _nh = _im.size
                    # background-origin:content-box paints inside the padding
                    # (the QR case); padding-box/border-box fill the whole box.
                    _origin = el.get("bg_origin") or "padding-box"
                    if "content" in _origin:
                        _pl = el.get("pad_left_px", 0); _pr = el.get("pad_right_px", 0)
                        _pt = el.get("pad_top_px", 0); _pb = el.get("pad_bottom_px", 0)
                    else:
                        _pl = _pr = _pt = _pb = 0
                    _cx = el["x"] + _pl; _cy = el["y"] + _pt
                    _cw = max(1, el["w"] - _pl - _pr)
                    _ch = max(1, el["h"] - _pt - _pb)
                    _size = (el.get("bg_size") or "auto").strip().lower()
                    _crop = None
                    pw, ph, ox, oy = _cw, _ch, 0, 0
                    if _nw > 0 and _nh > 0:
                        _nr = _nw / _nh
                        _br = _cw / _ch
                        if "contain" in _size:
                            if _nr > _br:
                                pw, ph = _cw, _cw / _nr
                            else:
                                pw, ph = _ch * _nr, _ch
                            ox = (_cw - pw) / 2; oy = (_ch - ph) / 2
                        elif "cover" in _size:
                            if _nr > _br:
                                _crop = ("lr", (1 - _br / _nr) / 2)
                            else:
                                _crop = ("tb", (1 - _nr / _br) / 2)
                    _px = emu_x(_cx + ox); _py = emu_y(_cy + oy)
                    _pw = emu_x(_cx + ox + pw) - _px
                    _ph = emu_y(_cy + oy + ph) - _py
                    _pic = slide.shapes.add_picture(_bgpath, _px, _py,
                                                    width=_pw, height=_ph)
                    if _crop:
                        _sd, _cv = _crop
                        if _sd == "lr":
                            _pic.crop_left = _cv; _pic.crop_right = _cv
                        else:
                            _pic.crop_top = _cv; _pic.crop_bottom = _cv
                    counts["pictures"] += 1
                except Exception as _e:
                    print(f"[skip bg-img {_bgpath}] {_e}", file=sys.stderr)

    # ─── Pass 2: emit TEXT BLOCKS (each block element = ONE TextBox with mixed runs) ──
    # text_blocks: list of {x, y, w, h, runs: [{text, font, ...}]}
    # Inline children (<strong>, <em>) become Runs in the same paragraph,
    # preserving inline-flow rendering.
    text_blocks = sorted(dom["text_blocks"], key=lambda t: (t["depth"], t["y"], t["x"]))
    align_map = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT,
                 "center": PP_ALIGN.CENTER, "justify": PP_ALIGN.JUSTIFY,
                 "start": PP_ALIGN.LEFT, "end": PP_ALIGN.RIGHT}

    # Parse CSS line-height computed value ("31.68px" / "normal") → ABSOLUTE Pt.
    # CSS line-height is multiplier of FONT-SIZE; PPT line_spacing multiplier
    # is of FONT'S NATURAL LINE HEIGHT (ascent+descent+linegap) — different.
    # Passing absolute Pt to python-pptx gives a 1:1 match with CSS.
    def _parse_lh_pt(lh_str: str, font_size_px: float) -> float | None:
        if not lh_str or lh_str == "normal":
            return None
        s = lh_str.strip()
        if s.endswith("px"):
            try:
                return float(s[:-2]) * 72.0 / 96.0  # px → pt
            except ValueError:
                return None
        try:
            # Unitless multiplier (rare in computed style, but handle)
            return float(s) * font_size_px * 72.0 / 96.0
        except ValueError:
            return None

    def _apply_text_transform(text: str, tt: str) -> str:
        if not tt or tt == "none":
            return text
        if tt == "uppercase":
            return text.upper()
        if tt == "lowercase":
            return text.lower()
        if tt == "capitalize":
            return text.title()
        return text

    # CSS `hyphens: auto` uses libhyphen to break long words at syllable
    # boundaries during wrap. PPT has no equivalent setting but DOES respect
    # OOXML soft hyphens (U+00AD): invisible character that becomes a visible
    # "-" only at wrap points. Pre-insert them with pyphen so PPT wraps
    # like the browser.
    try:
        import pyphen
        _hyph = pyphen.Pyphen(lang="en_US")
    except ImportError:
        _hyph = None

    def _hyphenate(text: str, min_word_len: int = 7) -> str:
        if not _hyph or not text:
            return text
        # Tokenize on whitespace, preserve interior punctuation in word
        out_words = []
        for w in text.split(" "):
            # Strip leading/trailing punctuation, hyphenate the core
            i, j = 0, len(w)
            while i < j and not w[i].isalnum():
                i += 1
            while j > i and not w[j-1].isalnum():
                j -= 1
            prefix, core, suffix = w[:i], w[i:j], w[j:]
            if len(core) >= min_word_len and core.isalpha():
                core = _hyph.inserted(core, "­")
            out_words.append(prefix + core + suffix)
        return " ".join(out_words)

    # Truncate text to fit content box height (for overflow:hidden / line-clamp).
    # Uses PIL with fc-match resolved font; falls back to char-count heuristic.
    def _truncate_to_box(text: str, font_name: str, font_size_pt: float,
                         bold: bool, content_w_px: float, content_h_px: float,
                         max_lines: int = 0) -> str:
        try:
            from PIL import ImageFont
            import subprocess
            spec = font_name + (":weight=700" if bold else "")
            out = subprocess.run(["fc-match", "-f", "%{file}", spec],
                                 capture_output=True, text=True, timeout=2)
            font_path = out.stdout.strip()
            if not font_path:
                return text
            font = ImageFont.truetype(font_path, max(6, int(round(font_size_pt))))
            line_h = font_size_pt * 1.30 * 96 / 72
            if max_lines > 0:
                allowed = max_lines
            else:
                allowed = max(1, int(content_h_px / line_h))
            # Word-wrap greedy
            words = text.split(" ")
            cur = ""
            lines = []
            for w in words:
                cand = (cur + " " + w).strip() if cur else w
                bb = font.getbbox(cand)
                if cur and (bb[2] - bb[0]) > content_w_px:
                    lines.append(cur)
                    cur = w
                    if len(lines) >= allowed:
                        break
                else:
                    cur = cand
            if len(lines) < allowed and cur:
                lines.append(cur)
            if len(lines) <= allowed:
                return text
            # Take allowed lines, add ellipsis at end
            kept = lines[:allowed]
            kept[-1] = kept[-1].rstrip() + "…"
            return " ".join(kept)
        except Exception:
            return text

    n_force_wrapped = 0
    for block_idx, tb_data in enumerate(text_blocks):
        font_scale = corrections.get(str(block_idx), {}).get("font_scale", 1.0)
        # Force PPT to wrap at exactly the browser's wrap positions by
        # splitting runs and inserting \n at each line boundary. Eliminates
        # wrap drift caused by font-shaper differences between browser and
        # PPT. Idempotent: marker `_wrap_forced` prevents re-splitting on
        # subsequent rounds (mutates tb_data in place across L2 iterations).
        # Headings are force-wrapped UNCONDITIONALLY (not gated on match_wrap):
        # the poster title (h1) is the worst overflow offender — PPT's wider font
        # advance re-wraps a 2-line browser title to 3 lines, and the extra line
        # spills down onto the author row. Honoring the browser's exact wrap
        # points keeps the title at its measured line count. Only fires when the
        # heading actually wrapped (>1 line), so single-line headings are untouched.
        # `.scan-link` etc.). A flex-column container's stacked spans are
        # visually separate lines the browser produced by LAYOUT, not by text
        # wrapping — PPT has no flex, so without explicit \n it collapses them
        # onto one line. Force-wrap such blocks too so extractWrapLines' per-
        # item Y splits become hard breaks.
        _force_wrap_this = (match_wrap
                            or tb_data.get("tag") in ("h1", "h2", "h3")
                            or tb_data.get("is_flex_col"))
        if _force_wrap_this and tb_data.get("wrap_lines") and len(tb_data["wrap_lines"]) > 1 \
                and not tb_data.get("_wrap_forced"):
            new_runs = _split_runs_at_browser_wraps(tb_data["runs"], tb_data["wrap_lines"])
            if new_runs is not tb_data["runs"]:
                tb_data["runs"] = new_runs
                tb_data["_wrap_forced"] = True
                n_force_wrapped += 1
        # Title-fit downscale: PPT/soffice render the same font wider than the
        # browser AND wrap with a different algorithm, so a multi-line h1 title
        # gains an extra line that overflows the fixed header band onto the author
        # row (worst on wide centered titles — v3/reel). Predicting the exact
        # foreign wrap is unreliable, so shrink a multi-line title by a fixed safety
        # factor scaled to how tightly it fills its widest line: the closer the
        # widest browser line is to the box width, the more shrink it needs to keep
        # its line count in PPT. Single-line titles and titles with lots of slack
        # are left ~untouched.
        if tb_data.get("tag") == "h1" and tb_data.get("wrap_lines") \
                and len([ln for ln in tb_data["wrap_lines"] if ln.strip()]) > 1:
            try:
                import subprocess as _sp
                from PIL import ImageFont as _IF
                _r0 = (tb_data.get("runs") or [{}])[0]
                _fs = float(_r0.get("font_size_px", 40) or 40)
                _fam = (_r0.get("font_family", "") or "Inter").split(",")[0].strip().strip('"\'')
                _bold = any(int(str(r.get("font_weight", 400)) or 400) >= 600
                            for r in (tb_data.get("runs") or []))
                _fp = _sp.run(["fc-match", "-f", "%{file}",
                               _fam + (":bold" if _bold else "")],
                              capture_output=True, text=True).stdout.strip()
                _bw = max(10.0, float(tb_data.get("w", 0))
                          - float(tb_data.get("pad_left_px", 0))
                          - float(tb_data.get("pad_right_px", 0)))
                if _fp:
                    _base = _IF.truetype(_fp, max(6, int(round(_fs))))
                    _maxlw = max((_base.getbbox(ln)[2] - _base.getbbox(ln)[0])
                                 for ln in tb_data["wrap_lines"] if ln.strip())
                    # The title is ALREADY hard-wrapped at the browser's exact
                    # line breaks (force-wrap above inserts \n), so PPT cannot
                    # gain an extra line by re-wrapping — each forced line is
                    # atomic. The only remaining risk is a single forced line
                    # being wider than the box under PPT's ~6% wider glyph
                    # advance, which would soft-wrap THAT line. So shrink only
                    # enough to keep the widest line inside the box (target 96%
                    # of box width, ×1.06 glyph fudge), NOT down to a low fill
                    # target — that was over-shrinking wide titles to barely
                    # above the author row (88pt→57pt on the reel poster).
                    _need = (_bw * 0.96) / max(1.0, _maxlw * 1.06)
                    if _need < 0.995:
                        font_scale *= max(0.6, _need)
            except Exception:
                pass
            except Exception:
                pass
        # For <li>: extend textbox LEFT by parent ul/ol's padding-left so
        # the native bullet can sit in the list's marker area, NOT at the
        # li content edge (which would align bullet with above text).
        list_shift_px = (tb_data.get("parent_list_pad_px", 0)
                         if tb_data.get("is_list_item") else 0)
        x = emu_x(tb_data["x"] - list_shift_px)
        y = emu_y(tb_data["y"])
        w = emu_x(tb_data["x"] + tb_data["w"]) - x
        h = emu_y(tb_data["y"] + tb_data["h"]) - y
        # Bold-aware width safety pad. Bold runs are the one place
        # font-metric drift is real: html_to_pptx remaps every weight
        # >= 600 (SemiBold / Bold / ExtraBold) onto the single 'Inter'
        # family with b='1', and font_embedder puts Inter-ExtraBold.ttf
        # in the bold slot — so a SemiBold (600) run in HTML gets
        # rendered at ExtraBold (800) advance in PPT, which is wider
        # than the bbox Chrome measured. Without compensation a tight
        # single-line bold heading ('faster sample ranking · COCO',
        # the Headline-Numbers hero label) wraps in PPT but didn't in
        # the browser. Regular-weight text (weight 400 → Inter Regular
        # in both) needs no pad — the old uniform pad pushed regular
        # text's edges off the HTML bbox for no benefit.
        #
        # Trigger is ANY styled run (reverted 2026-06-15 from ALL-only
        # — the ALL-only refinement under-padded paragraphs with just
        # one or two <strong>/<em> words and they still overflowed to
        # a new line in PPT. Extra horizontal slack is harmless;
        # overflowing the bbox into the next visible line is the actual
        # defect we need to prevent). Bold AND italic both drift —
        # Inter Italic glyph metrics differ from Inter Regular by enough
        # to push a tight line past the bbox edge, same way bold does.
        def _has_styled_run(rs):
            for r in rs:
                fw = r.get("font_weight")
                fs = (r.get("font_style") or "").lower()
                if fs in ("italic", "oblique"):
                    return True
                try:
                    if fw and int(fw) >= 600:
                        return True
                except (ValueError, TypeError):
                    if fw in ("bold", "bolder"):
                        return True
            return False

        def _has_non_embedded_font_run(rs):
            """True when any run uses a family we do NOT embed in the .pptx.

            Embedded families (see EMBEDDED_FONT_FAMILIES at module top)
            measure with the same binary that PowerPoint will render with,
            so glyph-advance is bit-identical and a tight single-line bbox
            stays tight. Non-embedded families use host fontconfig fallback
            for measurement (e.g. Liberation Serif standing in for Times
            New Roman on Linux) but render with the host OS's real font in
            Mac/Win PowerPoint — even nominally metric-compatible fallbacks
            drift 1-3% on individual glyphs, enough to push a tight label
            from "fits" to "overflows the textbox" or "wraps to two lines"
            after round-tripping through PPT.

            ANY-trigger semantics for the same reason as _has_styled_run:
            even one non-embedded run can push the line past the bbox edge
            and force a wrap. Update EMBEDDED_FONT_FAMILIES whenever the
            pipeline starts embedding a new family."""
            for r in rs:
                fam = (r.get("font_family") or "").lower()
                if not fam:
                    continue
                if not any(emb in fam for emb in EMBEDDED_FONT_FAMILIES):
                    return True
            return False

        # Compose width pad: take the LARGEST of all applicable mitigations.
        # - any bold OR italic run needs pad because we collapse 600/700/800
        #   onto the embedded ExtraBold slot, and italic has its own metric
        #   drift between browser measurement and PowerPoint render
        # - any non-embedded run needs pad because the browser-side fallback
        #   font has slightly different metrics from the Mac/Win PowerPoint
        #   font that ships with the host OS
        _pad_pct, _pad_cap_px = 0.0, 0
        if _has_styled_run(tb_data["runs"]):
            # Bumped 1.5%/18px → 2.5%/30px (2026-06-15) — at the previous
            # cap, lines mixing one Inter-Italic <em> with regular runs
            # still wrapped to a new line. Extra horizontal slack is
            # harmless per the user spec, overflow is the actual defect.
            _pad_pct, _pad_cap_px = max(_pad_pct, 0.025), max(_pad_cap_px, 30)
        if _has_non_embedded_font_run(tb_data["runs"]):
            # Bumped 1.8%/32px → 3.0%/48px (2026-06-15) — same rationale
            # as the styled-run pad; the older 32px cap left mixed long
            # author/affiliation strings still pushing past the bbox.
            _pad_pct, _pad_cap_px = max(_pad_pct, 0.030), max(_pad_cap_px, 48)
        _width_pad = (min(int(w * _pad_pct), emu_x(_pad_cap_px))
                      if _pad_cap_px > 0 else 0)
        w += _width_pad
        # Pad height a touch so descenders / wrapped lines don't clip
        h = int(h * 1.15)
        try:
            tb = slide.shapes.add_textbox(x, y, w, h)
            tf = tb.text_frame
            # Single-word short text in a tight bbox (e.g. "ICLR"/"ICML"
            # 4-letter badges at 88pt) — PPT auto-wrap kicks in when font
            # metric is slightly wider than browser, producing "ICL\nR".
            # Browser proved the word fits; force no-wrap so PPT honors it.
            single_word_short = False
            if tb_data.get("runs"):
                joined = "".join(r["text"] for r in tb_data["runs"] if r["text"] != "\n").strip()
                if joined and " " not in joined and len(joined) <= 6:
                    single_word_short = True
            # When we've ALREADY forced the wrap points (inserted synthetic \n
            # at the browser's exact line boundaries via
            # _split_runs_at_browser_wraps), PPT must NOT add its own wrapping on
            # top: with word_wrap=True a forced line that renders 1-3% wider in
            # PowerPoint's real font (vs the browser/soffice measurement) gets
            # re-wrapped, leaving one word stranded on the next line — the
            # "carriage return after the first word" staircase (reported on
            # latent-diffusion Method paragraphs). Same principle as the
            # single-short-word no-wrap below: the browser proved these lines
            # fit, so force-honor them. The width pad above keeps the slightly
            # wider PPT glyphs inside the box rather than spilling.
            force_no_wrap = single_word_short or bool(tb_data.get("_wrap_forced"))
            tf.word_wrap = not force_no_wrap
            # CSS padding → textbox internal margins (text inset from bbox edge)
            tf.margin_left = Emu(emu_x(tb_data.get("pad_left_px", 0)))
            tf.margin_right = Emu(emu_x(tb_data.get("pad_right_px", 0)))
            tf.margin_top = Emu(emu_y(tb_data.get("pad_top_px", 0)))
            tf.margin_bottom = Emu(emu_y(tb_data.get("pad_bottom_px", 0)))
            p_ = tf.paragraphs[0]
            ta = (tb_data.get("text_align") or "").lower()
            p_.alignment = align_map.get(ta, PP_ALIGN.LEFT)

            # For <li>: add OOXML native bullet (<a:buChar>) and hanging
            # indent so wrap-lines align with first-line text, while the
            # bullet sits in the freshly-extended-left area.
            if tb_data.get("is_list_item") and list_shift_px > 0:
                shift_emu = emu_x(list_shift_px)
                pPr = p_._pPr if p_._pPr is not None else p_._p.get_or_add_pPr()
                pPr.set("marL", str(shift_emu))
                pPr.set("indent", str(-shift_emu))
                # Use the section color if available (from the li::marker
                # in CSS would require a per-section lookup; just use the
                # first run's color as a reasonable proxy)
                first_run = tb_data["runs"][0] if tb_data["runs"] else None
                marker_color = None
                if first_run:
                    c = _parse_color(first_run.get("color") or "")
                    if c:
                        marker_color = c[0]
                # Remove any existing bullet config, then add ours
                for child in list(pPr):
                    if child.tag in (qn("a:buChar"), qn("a:buFont"),
                                     qn("a:buNone"), qn("a:buClr")):
                        pPr.remove(child)
                if marker_color:
                    buClr = etree.SubElement(pPr, qn("a:buClr"))
                    etree.SubElement(buClr, qn("a:srgbClr"),
                                     val=f"{marker_color}")
                etree.SubElement(pPr, qn("a:buFont"), typeface="Arial",
                                 panose="020B0604020202020204",
                                 pitchFamily="34", charset="0")
                etree.SubElement(pPr, qn("a:buChar"), char="•")

            # CSS line-height → PPT paragraph.line_spacing absolute Pt.
            # Browser bbox was measured at CSS line-height (e.g., 1.36 × 24pt
            # = 32.64pt per line). PPT default = font's natural line height
            # (Inter ≈ 1.21 × 24pt = 29pt). The ~3-4pt-per-line shortfall
            # accumulates: each textbox's text under-fills its bbox, leaving
            # trailing whitespace; subsequent blocks sit at their browser-bbox
            # positions (lower), so visible paragraph gaps look 10-20px wider
            # than HTML.
            #
            # Setting line_spacing = Pt(css_line_height_in_pt) makes the PPT
            # baseline-to-baseline distance match CSS, eliminating the
            # trailing whitespace. Earlier this looked worse — but only
            # because the font was DejaVu Sans (wider per char, wrapping to
            # extra lines, overflowing). With Inter installed it works.
            first_run = tb_data["runs"][0] if tb_data["runs"] else None
            lh_pt = _parse_lh_pt(tb_data.get("line_height") or "",
                                 first_run["font_size_px"] if first_run else 16.0)
            if lh_pt:
                # Apply font_scale to line spacing too — shrunk text needs
                # proportionally shrunk line spacing or the box still under-fills
                p_.line_spacing = Pt(lh_pt * font_scale * slide_scale)
            # Force 0 — PPT layout default is non-zero; we want browser-bbox
            # gaps to be the only source of vertical spacing between blocks.
            p_.space_before = Pt(0)
            p_.space_after = Pt(0)
            tt = tb_data.get("text_transform", "none")
            hyphens_auto = tb_data.get("hyphens") == "auto"

            # For heading text blocks containing decoration runs (e.g., span.num
            # giving "1" in a different color from the heading), force a
            # consistent color (h3's own text color) so the decorative number
            # doesn't disappear against the banner bg. Also separate runs with
            # a space when they come from element boundaries so "1Motivation"
            # becomes "1 Motivation".
            is_heading = tb_data["tag"] in ("h1", "h2", "h3", "h4", "h5", "h6")
            heading_color = None
            if is_heading and len(tb_data["runs"]) >= 2:
                # The h3's own text color is in the LAST text-node run (text
                # nodes directly under h3 use h3's color; nested span text
                # uses span's color).
                heading_color = tb_data["runs"][-1].get("color")
                # Inject a space between runs to avoid "1Motivation"
                new_runs = []
                for i, r in enumerate(tb_data["runs"]):
                    if i > 0 and r["text"] != "\n":
                        prev = new_runs[-1]["text"] if new_runs else ""
                        if prev and not prev.endswith((" ", "\n", " ")) \
                                and not r["text"].startswith((" ", "\n", " ")):
                            sep = dict(r)
                            sep["text"] = " "
                            new_runs.append(sep)
                    new_runs.append(r)
                tb_data["runs"] = new_runs

            # If overflow:hidden / -webkit-line-clamp: truncate the joined
            # text content to what would fit, then add "…" — matches what
            # the browser shows (clamped + ellipsis) rather than pasting
            # the full hidden text and overflowing the bbox.
            overflow_hidden = tb_data.get("overflow_hidden", False)
            line_clamp = tb_data.get("webkit_line_clamp", 0)
            if (overflow_hidden or line_clamp > 0) and tb_data["runs"]:
                first_run = tb_data["runs"][0]
                # Join all runs into one string for measurement
                joined = "".join(r["text"] for r in tb_data["runs"] if r["text"] != "\n")
                font_name = (first_run.get("font_family") or "").split(",")[0].strip().strip('"').strip("'")
                fs_px = first_run.get("font_size_px", 16) * font_scale
                fs_pt = fs_px * 72 / 96
                bold = False
                fw = first_run.get("font_weight", "400")
                try:
                    bold = int(fw) >= 600
                except (ValueError, TypeError):
                    bold = fw in ("bold", "bolder")
                content_w_px = max(10, tb_data["w"]
                                   - tb_data.get("pad_left_px", 0)
                                   - tb_data.get("pad_right_px", 0))
                content_h_px = max(10, tb_data["h"]
                                   - tb_data.get("pad_top_px", 0)
                                   - tb_data.get("pad_bottom_px", 0))
                truncated = _truncate_to_box(
                    joined, font_name, fs_pt, bold,
                    content_w_px, content_h_px, max_lines=line_clamp)
                if truncated != joined:
                    # Replace with a single truncated run (lose intra-block
                    # styling — acceptable for clamped captions which are
                    # typically single-style anyway).
                    tb_data["runs"] = [{
                        "text": truncated,
                        "font_family": first_run.get("font_family", ""),
                        "font_size_px": first_run.get("font_size_px", 16),
                        "font_weight": fw,
                        "font_style": first_run.get("font_style", "normal"),
                        "color": first_run.get("color", "rgb(0,0,0)"),
                        "letter_spacing_px": first_run.get("letter_spacing_px", 0),
                    }]

            # CSS ::marker (bullet) is generated content — not in textContent.
            # For <li>: use OOXML native bullet via buChar, extend textbox LEFT
            # by parent ul/ol's padding-left so bullet sits in the list's
            # marker area (matches HTML where ::marker sits in ul's padding,
            # not in li content area). Hanging-indent keeps wrap-line text
            # aligned with first line.
            if tb_data.get("is_list_item") and tb_data["runs"]:
                # tb already positioned at li.bbox.x; need to redo with shifted x
                pass  # handled below — sentinel for clarity

            # Each "run" is a (text, style) chunk. Add runs in order to the
            # SAME paragraph so they flow inline like in HTML.
            for run_data in tb_data["runs"]:
                if run_data["text"] == "\n":
                    # <br> → new paragraph (PPT doesn't have inline <br>)
                    p_ = tf.add_paragraph()
                    p_.alignment = align_map.get(ta, PP_ALIGN.LEFT)
                    if lh_pt:
                        p_.line_spacing = Pt(lh_pt * font_scale * slide_scale)
                    p_.space_before = Pt(0)
                    p_.space_after = Pt(0)
                    # For bulleted list items: the FIRST paragraph above got
                    # marL + negative indent + buChar so the bullet sits in
                    # the extended-left area and first-line text starts at
                    # marL. Continuation paragraphs (from browser word-wrap)
                    # must inherit the SAME marL so their text aligns with
                    # the first-line text, not with the bullet position —
                    # AND must suppress the bullet (<a:buNone/>) so PPT
                    # doesn't repeat "•" on each wrapped line. Without this
                    # the wrapped lines start at the textbox edge, breaking
                    # the hanging-indent visual that matches HTML's <ul><li>
                    # rendering (user-reported 2026-06-11).
                    if tb_data.get("is_list_item") and list_shift_px > 0:
                        pPr_c = (p_._pPr if p_._pPr is not None
                                 else p_._p.get_or_add_pPr())
                        pPr_c.set("marL", str(emu_x(list_shift_px)))
                        pPr_c.set("indent", "0")
                        for child in list(pPr_c):
                            if child.tag in (qn("a:buChar"), qn("a:buFont"),
                                             qn("a:buNone"), qn("a:buClr")):
                                pPr_c.remove(child)
                        etree.SubElement(pPr_c, qn("a:buNone"))
                    continue
                run = p_.add_run()
                # Per-run text_transform override (e.g. ::before { text-transform:
                # uppercase } on a callout prefix inside a non-uppercase paragraph).
                run_tt = run_data.get("text_transform") or tt
                run.text = _apply_text_transform(run_data["text"], run_tt)
                if hyphens_auto:
                    run.text = _hyphenate(run.text)
                font = run.font
                ff = run_data.get("font_family") or ""
                first_family = _pick_font_family(ff)
                # OOXML's classic embedded-font spec gives each typeface only
                # 4 slots (Regular / Bold / Italic / BoldItalic) — there's
                # no slot for SemiBold (600) or ExtraBold (800). If we just
                # rely on font.bold=True for weight>=600, both SemiBold (600)
                # and Bold (700) collapse to the same embedded Bold .ttf,
                # losing the visual weight differentiation AND making
                # SemiBold-styled text render with Bold's slightly wider
                # x-advance (which then wraps to extra lines in PPTX vs HTML).
                # Workaround: when the run's family is Inter, remap weight
                # 600 -> typeface "Inter SemiBold" and weight 800 -> typeface
                # "Inter ExtraBold". These are embedded as SEPARATE typeface
                # families (each occupying its own family's Regular slot)
                # via font_embedder's repeatable --font flag — see the
                # html2pptx Inter-embed step. Bold (700) still
                # uses the regular "Inter" family with bold=True.
                #
                # IMPORTANT: Chrome's getComputedStyle reports the matched
                # PostScript font name, not the CSS family. On Linux where
                # Inter is installed as separate per-weight files (e.g.
                # /usr/share/fonts/.../Inter-ExtraBold.ttf), Chrome reports
                # 'Inter-ExtraBold' for ALL runs styled `font-family: Inter`
                # (because it picks one available weight and returns that
                # font's PostScript name). So before checking, normalize
                # the family by stripping known weight suffixes — get back
                # to the base family ('Inter') and then apply the
                # weight-based remap from the numeric font-weight CSS value.
                base_family = first_family
                if base_family:
                    for _suffix in ("-ExtraBold", "-Extra Bold", "-SemiBold",
                                    "-Semi Bold", "-Bold", "-Regular",
                                    "-Medium", "-Light", "-Thin", "-Black",
                                    " ExtraBold", " SemiBold", " Bold",
                                    " Regular", " Medium"):
                        if base_family.endswith(_suffix):
                            base_family = base_family[: -len(_suffix)]
                            break
                fw_raw = run_data.get("font_weight")
                try:
                    fw_int = int(fw_raw) if fw_raw else None
                except (ValueError, TypeError):
                    fw_int = None
                # Single-family + bold-flag strategy. We previously tried
                # multi-typeface embedding ('Inter SemiBold' / 'Inter
                # ExtraBold' as separate families) so PPTX would 1:1 mirror
                # HTML's weight ladder, but PowerPoint and LibreOffice both
                # silently fail to load alt-family embedded fonts when the
                # OS doesn't already have those families installed —
                # OOXML's <p:embeddedFontLst> matching is brittle.
                #
                # The pipeline now embeds a single 'Inter' typeface with
                # the bold slot pointing at Inter-ExtraBold.ttf (see the
                # html2pptx Inter-embed step). Any weight >= 600 here
                # becomes typeface='Inter' + b='1', which makes PPT load
                # the bold slot — guaranteed to render visibly heavier
                # than regular regardless of whether the embedded font
                # was successfully loaded by the host (because the host
                # can always fall back to system Inter Bold or synthesize
                # bold). We lose the 600/700/800 visual ladder, but
                # 'is this run bold' is preserved 1:1 with the HTML.
                if base_family and base_family.lower() == "inter":
                    first_family = "Inter"
                if first_family:
                    font.name = first_family
                if run_data.get("font_size_px"):
                    # Apply both block-level font_scale (L2 corrections) AND
                    # global slide-scale (poster cap from 60→55 etc.).
                    fsz_px = run_data["font_size_px"]
                    # Sup/sub readability floor: CSS often sets <sup> to
                    # 0.65em which becomes very tiny on a poster. Combined
                    # with the 30% baseline raise, glyph effective size is
                    # ~0.45em on screen — illegible at print scale.
                    # Floor at 1.4× parent so post-shift visual is roughly
                    # parent size (10pt parent → 14pt sup → raised 30%
                    # gives an effective glyph footprint near 10pt). Per
                    # user preference on academic posters where every mark
                    # must read at distance.
                    va = run_data.get("vertical_align", "baseline")
                    if va in ("super", "sub"):
                        parent_sz = None
                        for prev in tb_data["runs"]:
                            if prev is run_data:
                                break
                            if (prev.get("text") != "\n"
                                    and prev.get("vertical_align", "baseline") == "baseline"
                                    and prev.get("font_size_px")):
                                parent_sz = prev["font_size_px"]
                        if parent_sz is None:
                            for r0 in tb_data["runs"]:
                                if r0.get("text") != "\n" and r0.get("font_size_px"):
                                    parent_sz = r0["font_size_px"]
                                    break
                        if parent_sz and fsz_px < parent_sz * 1.4:
                            fsz_px = parent_sz * 1.4
                    font.size = Pt(fsz_px * 72 / 96 * font_scale * slide_scale)
                fw = run_data.get("font_weight")
                if fw:
                    try:
                        font.bold = int(fw) >= 600
                    except (ValueError, TypeError):
                        font.bold = fw in ("bold", "bolder")
                if run_data.get("font_style") == "italic":
                    font.italic = True
                color = _parse_color(run_data.get("color") or "")
                if color:
                    font.color.rgb = color[0]
                # For headings: override per-run color with the heading's own
                # color so decorative runs (span.num "1") aren't invisible
                # against the banner background.
                if is_heading and heading_color:
                    hc = _parse_color(heading_color)
                    if hc:
                        font.color.rgb = hc[0]
                # CSS letter-spacing → OOXML run rPr spc
                ls = run_data.get("letter_spacing_px", 0)
                if ls and abs(ls) >= 0.05:
                    _set_run_letter_spacing(run, ls * slide_scale, run_data.get("font_size_px", 16))
                # Superscript / subscript: CSS vertical-align: super|sub OR
                # <sup>/<sub> element → PPT run baseline offset.
                va = run_data.get("vertical_align", "baseline")
                if va in ("super", "sub"):
                    _set_run_baseline(run, va)
                # Inline highlight: CSS background-color on an inline element
                # → PPT run highlight (NOT a standalone background rectangle).
                bgc = _parse_color(run_data.get("bg_color") or "")
                if bgc:
                    _set_run_highlight(run, bgc[0])
            counts["text"] += 1
        except Exception as e:
            print(f"[skip text block {tb_data['tag']}] {e}", file=sys.stderr)

    # ─── Pass 3: emit MATH BLOCKS as native PowerPoint OMML equations ────
    # The JS DOM extractor populated dom['math_blocks'] from every
    # <mjx-container data-tex="..." data-display="0|1"> in the rendered
    # page. Convert each TeX string → MathML (latex2mathml) → OMML
    # (mathml2omml), wrap in <a14:m> (DrawingML 2010 math extension —
    # PowerPoint slides ONLY recognize OMML inside this wrapper, not as
    # a bare <m:oMath> paragraph child; tested 2026-06-15 with E=mc²
    # round-trip), and inject as a textbox at the captured bbox. The
    # equation is then editable in PowerPoint via Insert→Equation.
    #
    # Graceful degradation: if latex2mathml / mathml2omml aren't
    # installed, we log once and skip — the equation will already exist
    # in the rendered PNG (since MathJax SVG is rasterized elsewhere),
    # the .pptx just loses the native math affordance.
    math_blocks = dom.get("math_blocks") or []
    if math_blocks and not _MATH_AVAILABLE:
        # No native-OMML capability. We can still salvage any equation that
        # was rasterized during extraction (mb["png"]); only those without a
        # raster are truly lost, so warn but DO NOT clear the whole list.
        _no_raster = sum(1 for m in math_blocks if not m.get("png"))
        if _no_raster:
            print(f"[math skip] {_no_raster}/{len(math_blocks)} equation(s) "
                  "found but latex2mathml / mathml2omml not installed and no "
                  "raster fallback — `pip install latex2mathml mathml2omml` "
                  "for native PPT math", file=sys.stderr)
    for mb in math_blocks:
        # Universal fallback: if extraction captured a raster of the math
        # region, emit it as a Picture. Unlike native OMML (<a14:m>, which
        # only Microsoft PowerPoint renders — soffice/Keynote/Google Slides
        # show it blank), a PNG renders everywhere. Prefer it so the exported
        # PNG / non-PowerPoint viewers are never blank. Falls through to OMML
        # only when no raster exists.
        _mpng = mb.get("png")
        if _mpng:
            try:
                _mx = emu_x(mb["x"]); _my = emu_y(mb["y"])
                _mw = max(1, emu_x(mb["x"] + mb["w"]) - _mx)
                _mh = max(1, emu_y(mb["y"] + mb["h"]) - _my)
                slide.shapes.add_picture(_mpng, _mx, _my, width=_mw, height=_mh)
                continue
            except Exception as e:
                print(f"[math raster fail] {mb.get('tex','')[:40]!r}: {e}",
                      file=sys.stderr)
        if not _MATH_AVAILABLE:
            continue
        try:
            mml = _l2m.convert(mb["tex"],
                               display="block" if mb["display"] else "inline")
            omath_xml = _m2o.convert(mml)
        except Exception as e:
            print(f"[math skip] tex={mb['tex'][:60]!r}: {e}", file=sys.stderr)
            continue
        try:
            # Size: PowerPoint renders OMML at the application default math
            # font size (~10pt) unless each <m:r> carries an explicit
            # rPr/sz. The mjx-container's getComputedStyle font-size (the
            # author's intent, e.g. 1.25em of section body ≈ 55pt at this
            # canvas) lives in mb["font_size_px"]. Inject sz="<half-pt
            # ×100>" into every m:r so the equation renders at that size,
            # not tiny.  slide_scale already shrunk by the OOXML-clamp
            # factor; apply it here too so the math matches the rest of
            # the slide's typography.
            fsz_px = mb.get("font_size_px") or 16
            sz_pt = fsz_px * 72.0 / 96.0 * slide_scale
            sz_attr = max(100, int(round(sz_pt * 100)))  # OOXML sz = half-pt ×100
            # Inject into every <m:r> that doesn't already have a sz attr.
            # Use a non-greedy regex so we don't disturb existing rPr.
            def _add_sz(m: re.Match) -> str:
                pre, body, post = m.group(1), m.group(2), m.group(3)
                if 'sz="' in body:
                    return m.group(0)
                # Insert <m:rPr> ... </m:rPr> with sz if no rPr; otherwise
                # add sz to the existing rPr's first <a:rPr>.
                if "<m:rPr>" in body:
                    body2 = body.replace(
                        "<m:rPr>",
                        f'<m:rPr><a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" sz="{sz_attr}"/>',
                        1,
                    )
                else:
                    body2 = (
                        f'<m:rPr><a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" sz="{sz_attr}"/></m:rPr>'
                        + body
                    )
                return f'{pre}{body2}{post}'
            omath_xml = re.sub(
                r'(<m:r>)(.*?)(</m:r>)', _add_sz, omath_xml, flags=re.DOTALL
            )

            x = emu_x(mb["x"]); y = emu_y(mb["y"])
            w = emu_x(mb["x"] + mb["w"]) - x
            h = emu_y(mb["y"] + mb["h"]) - y
            tb = slide.shapes.add_textbox(x, y, max(w, Emu(1)), max(h, Emu(1)))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.margin_left = Emu(0); tf.margin_right = Emu(0)
            tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]
            # Display math → <m:oMathPara> (own paragraph), inline → bare <m:oMath>
            inner = (f'<m:oMathPara xmlns:m="{_M_NS}">{omath_xml}</m:oMathPara>'
                     if mb["display"] else omath_xml)
            a14_wrap = (f'<a14:m xmlns:a14="{_A14_NS}" '
                        f'xmlns:m="{_M_NS}">{inner}</a14:m>')
            a14m = _lxml_etree.fromstring(a14_wrap)
            p._p.append(a14m)
            counts["math"] = counts.get("math", 0) + 1
        except Exception as e:
            print(f"[math emit fail] tex={mb['tex'][:60]!r}: {e}",
                  file=sys.stderr)

    prs.save(str(out_path))
    counts["force_wrapped"] = n_force_wrapped
    return counts


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--html", required=True, type=Path)
    ap.add_argument("--out", required=True, type=Path)
    ap.add_argument("--width-inch", type=float, default=None,
                    help="Slide width in inches. Auto-detect from CSS @page size if omitted.")
    ap.add_argument("--height-inch", type=float, default=None,
                    help="Slide height in inches. Auto-detect from CSS @page size if omitted.")
    ap.add_argument("--corrections", type=Path, default=None,
                    help="JSON file of L2 per-block corrections")
    ap.add_argument("--dom-cache", type=Path, default=None,
                    help="JSON cache of extracted DOM (skip browser if exists)")
    ap.add_argument("--match-wrap", action="store_true",
                    help="Force PPT to wrap text at the same character positions "
                         "Chrome wrapped at (uses extractWrapLines DOM data + "
                         "_split_runs_at_browser_wraps). Without this, the browser "
                         "and PowerPoint font shapers give slightly different glyph "
                         "widths, so PPT fits more chars per line — collapsing the "
                         "rendered height of multi-line callouts/paragraphs and "
                         "distorting the relative size proportions across the poster.")
    a = ap.parse_args()

    # Resolve canvas dimensions: explicit CLI > CSS @page > A0 default
    w_inch, h_inch = a.width_inch, a.height_inch
    if w_inch is None or h_inch is None:
        detected = detect_canvas_size(a.html)
        if detected:
            print(f"[canvas] auto-detected from @page: {detected[0]:.1f}×{detected[1]:.1f}\"",
                  file=sys.stderr)
            w_inch = w_inch or detected[0]
            h_inch = h_inch or detected[1]
        else:
            print(f"[canvas] no @page in HTML — falling back to A0 {DEFAULT_W_INCH}×{DEFAULT_H_INCH}\"",
                  file=sys.stderr)
            w_inch = w_inch or DEFAULT_W_INCH
            h_inch = h_inch or DEFAULT_H_INCH

    # OOXML hard-caps each slide dimension at 56 inches (51,206,400 EMU).
    # Conference posters routinely exceed this (60×36 in NeurIPS / CVPR;
    # 33.1×46.8 in A0 portrait fits). When either axis is over the cap,
    # uniformly down-scale the SLIDE to the tightest fit so the slide
    # stays inside OOXML limits AND the design aspect ratio is preserved.
    # The build_pptx slide_scale machinery (font/border/line auto-scaling)
    # then compensates downstream so the rendered poster looks the same
    # — it's just a smaller paper canvas.
    #
    # CRITICAL: the Playwright viewport MUST stay at the original design
    # size (design_w_inch / design_h_inch), NOT the clamped slide size.
    # If we shrink the viewport too, the browser renders the page into a
    # smaller box than its CSS @page expects → either the body overflows
    # OR text appears proportionally larger relative to the viewport, AND
    # build_pptx's `slide_scale = width_inch / (viewport_w/96)` ends up
    # = 1.0 (no scaling) instead of the intended ratio. Result: PPT body
    # text renders at the full designed pt size on a smaller slide →
    # visibly cramped sections (the user-reported "Problem section cramped" bug,
    # 2026-06-15). Pass design_* to extract_dom, slide_w/h_inch to build.
    OOXML_MAX_IN = PPT_MAX_INCH  # single source of truth (== cap_to_pptx_slide)
    design_w_inch, design_h_inch = w_inch, h_inch
    if w_inch > OOXML_MAX_IN or h_inch > OOXML_MAX_IN:
        scale = min(OOXML_MAX_IN / w_inch, OOXML_MAX_IN / h_inch)
        new_w, new_h = w_inch * scale, h_inch * scale
        print(f"[canvas] clamping slide {w_inch:.1f}×{h_inch:.1f}\" → "
              f"{new_w:.2f}×{new_h:.2f}\" (OOXML cap = {OOXML_MAX_IN}\"; "
              f"uniform scale {scale:.4f}×); design viewport kept at "
              f"{design_w_inch:.1f}×{design_h_inch:.1f}\"",
              file=sys.stderr)
        w_inch, h_inch = new_w, new_h

    if a.dom_cache and a.dom_cache.exists():
        print(f"[1/2] Loading DOM cache from {a.dom_cache}...", file=sys.stderr)
        dom = json.loads(a.dom_cache.read_text())
    else:
        print(f"[1/2] Extracting DOM from {a.html}...", file=sys.stderr)
        # Render at DESIGN size so CSS pt values (44pt body, 92pt titles)
        # land at their intended fraction of the canvas. build_pptx will
        # then scale down to fit the clamped slide via slide_scale.
        dom = extract_dom(a.html,
                          viewport_w=int(design_w_inch * 96),
                          viewport_h=int(design_h_inch * 96))
        if a.dom_cache:
            a.dom_cache.write_text(json.dumps(dom))
    print(f"      {len(dom['elements'])} visible elements, "
          f"body {dom['body_w']}×{dom['body_h']}px", file=sys.stderr)

    corrections = json.loads(a.corrections.read_text()) if a.corrections and a.corrections.exists() else {}
    if corrections:
        print(f"      applying {len(corrections)} block corrections", file=sys.stderr)

    print(f"[2/2] Building PPTX...", file=sys.stderr)
    counts = build_pptx(dom, a.out, w_inch, h_inch, corrections, match_wrap=a.match_wrap)

    print(f"\nWrote {a.out}", file=sys.stderr)
    parts = [f"{n} {k}" for k, n in counts.items() if n]
    print(f"  {', '.join(parts)}", file=sys.stderr)


if __name__ == "__main__":
    main()
