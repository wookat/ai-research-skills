#!/usr/bin/env python3
"""
Extract a slide-element registry from the final PPTX.

This registry is used by paper2video visual-cue planning. SVG group IDs are
useful semantic hints, and PPTX shape geometry gives the editable deck a second
auditable surface for where a highlight can safely land.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

SCHEMA_VERSION = "paper2video_pptx_elements.v1"
ANCHOR_RE = re.compile(r"\bcue_[A-Za-z0-9][A-Za-z0-9_.:-]*")


def _require_pptx():
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:  # pragma: no cover - depends on local env
        raise SystemExit(
            "[extract_pptx_elements] python-pptx is not installed in this Python env. "
            "Install python-pptx or run inside the ACL26 environment."
        ) from exc
    return Presentation


def clean_text(raw: str) -> str:
    return " ".join((raw or "").split())


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def shape_nonvisual_props(shape: Any) -> dict[str, str]:
    """Read cNvPr metadata used for shape names and alt text.

    PowerPoint stores shape label metadata under the non-visual property
    element, usually as `name`, `title`, and `descr` attributes. python-pptx
    exposes `shape.name`, but not every alt-text field is surfaced directly, so
    we read the underlying XML too.
    """
    props: dict[str, str] = {}
    element = getattr(shape, "element", None)
    if element is None:
        return props
    try:
        iterator = element.iter()
    except Exception:
        return props
    for node in iterator:
        if local_name(str(getattr(node, "tag", ""))) != "cNvPr":
            continue
        for key in ("name", "title", "descr"):
            value = clean_text(str(node.get(key) or ""))
            if value:
                props[key] = value
        break
    return props


def cue_anchors(*values: str) -> list[str]:
    anchors: set[str] = set()
    for value in values:
        anchors.update(ANCHOR_RE.findall(value or ""))
    return sorted(anchors)


def clamp01(value: float) -> float:
    return max(0.0, min(1.0, value))


def shape_type_name(shape: Any) -> str:
    try:
        return shape.shape_type.name
    except Exception:
        return str(getattr(shape, "shape_type", "UNKNOWN"))


def infer_role(name: str, text: str, shape_type: str, box: list[float]) -> str:
    x, y, w, h = box
    area = w * h
    blob = f"{name} {text} {shape_type}".lower()

    if area > 0.86 and not text:
        return "background"
    if y > 0.91 or re.search(r"\b\d{1,2}\s*/\s*\d{1,2}\b", text):
        return "footer"
    if text.lower().startswith("source:"):
        return "footer"
    if y < 0.21 and text and "TEXT" in shape_type:
        return "header"
    if area < 0.035 and re.search(r"\b(?:figure|table)\s+\d+", text, flags=re.IGNORECASE):
        return "caption"
    if any(k in blob for k in ("qr", "github", "arxiv", "repository", "repo", "code link", "paper link")):
        return "qr"
    if any(k in blob for k in ("formula", "equation", "sinusoid", "sine", "cosine", "wavelength", "positional encoding")):
        return "formula"
    role_map = [
        ("figure", ("figure", "image", "picture", "chart", "plot", "scaling", "evidence", "trajectory")),
        ("result", ("kpi", "metric", "result", "headline", "number", "accuracy", "average", "gain")),
        ("method", ("method", "pipeline", "flow", "step", "score", "select", "selection", "organize", "stream")),
        ("guidance", ("guidance", "boundary", "cyclic", "continuity", "diversity", "framework", "rule", "jitter", "jit", "local diversity", "homogeneous", "similar", "strictly sorted", "brittle")),
        ("takeaway", ("takeaway", "closing", "core", "claim", "toolbox")),
        ("title", ("title", "cover", "acl arr", "submission")),
    ]
    for role, keys in role_map:
        if any(k in blob for k in keys):
            return role
    if "TEXT" in shape_type and text:
        return "text"
    if "PICTURE" in shape_type:
        return "figure"
    return "content"


def read_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return clean_text(shape.text)
    except Exception:
        return ""


def shape_box(shape: Any, slide_w: int, slide_h: int) -> list[float] | None:
    try:
        left = float(shape.left)
        top = float(shape.top)
        width = float(shape.width)
        height = float(shape.height)
    except Exception:
        return None
    if width <= 0 or height <= 0 or slide_w <= 0 or slide_h <= 0:
        return None
    x = clamp01(left / slide_w)
    y = clamp01(top / slide_h)
    w = clamp01(width / slide_w)
    h = clamp01(height / slide_h)
    if w <= 0 or h <= 0:
        return None
    return [round(x, 6), round(y, 6), round(w, 6), round(h, 6)]


def should_keep_element(element: dict[str, Any]) -> bool:
    box = element["box"]
    area = box[2] * box[3]
    role = element["role"]
    text = element.get("text", "")
    shape_type = element.get("shape_type", "")

    if element.get("cue_anchors"):
        return True
    if role == "background":
        return False
    if area < 0.00045 and not text:
        return False
    if shape_type == "GROUP" and area > 0.84 and not text:
        return False
    return True


def walk_shapes(
    shapes: Any,
    *,
    slide_index: int,
    slide_w: int,
    slide_h: int,
    parent_id: str | None = None,
    z_prefix: tuple[int, ...] = (),
    depth: int = 0,
) -> list[dict[str, Any]]:
    elements: list[dict[str, Any]] = []
    for z, shape in enumerate(shapes):
        z_path = (*z_prefix, z)
        box = shape_box(shape, slide_w, slide_h)
        if box is None:
            continue

        sid = f"s{slide_index:02d}_sh{getattr(shape, 'shape_id', z)}"
        if parent_id:
            sid = f"{parent_id}_{getattr(shape, 'shape_id', z)}"
        name = clean_text(str(getattr(shape, "name", sid)))
        nv_props = shape_nonvisual_props(shape)
        if nv_props.get("name"):
            name = nv_props["name"]
        shape_type = shape_type_name(shape)
        text = read_text(shape)
        alt_title = nv_props.get("title", "")
        alt_description = nv_props.get("descr", "")
        metadata_text = clean_text(" ".join([name, alt_title, alt_description, text]))
        element = {
            "id": sid,
            "shape_id": int(getattr(shape, "shape_id", z)),
            "name": name,
            "alt_title": alt_title,
            "alt_description": alt_description,
            "cue_anchors": cue_anchors(name, alt_title, alt_description, text),
            "shape_type": shape_type,
            "parent_id": parent_id,
            "depth": depth,
            "z_order": list(z_path),
            "box": box,
            "point": [round(box[0] + box[2] / 2.0, 6), round(box[1] + box[3] / 2.0, 6)],
            "text": text[:500],
            "semantic_text": metadata_text[:1000],
            "role": infer_role(name, metadata_text, shape_type, box),
        }
        if should_keep_element(element):
            elements.append(element)

        if hasattr(shape, "shapes"):
            elements.extend(
                walk_shapes(
                    shape.shapes,
                    slide_index=slide_index,
                    slide_w=slide_w,
                    slide_h=slide_h,
                    parent_id=sid,
                    z_prefix=z_path,
                    depth=depth + 1,
                )
            )
    return elements


def extract_pptx_elements(pptx_path: Path) -> dict[str, Any]:
    Presentation = _require_pptx()
    prs = Presentation(str(pptx_path))
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    slides = []
    for idx in range(len(prs.slides)):
        slide = prs.slides[idx]
        elements = walk_shapes(
            slide.shapes,
            slide_index=idx + 1,
            slide_w=slide_w,
            slide_h=slide_h,
        )
        slides.append({
            "index": idx + 1,
            "id": f"slide_{idx + 1:02d}",
            "element_count": len(elements),
            "elements": elements,
        })
    return {
        "schema_version": SCHEMA_VERSION,
        "pptx": str(pptx_path),
        "slide_size_emu": [slide_w, slide_h],
        "slides": slides,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract PPTX slide element geometry for paper2video.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--out", type=Path, required=True)
    args = parser.parse_args()

    pptx = args.pptx.resolve()
    if not pptx.is_file():
        sys.exit(f"[extract_pptx_elements] PPTX not found: {pptx}")
    payload = extract_pptx_elements(pptx)
    args.out.parent.mkdir(parents=True, exist_ok=True)
    args.out.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(f"[extract_pptx_elements] wrote {sum(s['element_count'] for s in payload['slides'])} elements to {args.out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
