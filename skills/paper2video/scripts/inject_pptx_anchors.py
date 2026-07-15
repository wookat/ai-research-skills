#!/usr/bin/env python3
"""Inject visual-cue anchors into a final PPTX deck.

The highlighted video is rendered from PPTX, not directly from SVG. SVG group
ids are useful during authoring, but final QA must be able to resolve each
semantic cue to PPTX geometry. This script bridges the two surfaces by adding
invisible PPTX shapes at the cue-plan target boxes and writing the cue
`anchor_id` into standard PowerPoint non-visual metadata.
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any


SCHEMA_VERSION = "paper2video_pptx_anchor_injection.v1"
ANCHOR_RE = re.compile(r"\bcue_[A-Za-z0-9][A-Za-z0-9_.:-]*")


def utc_now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def clean_text(raw: str, limit: int = 260) -> str:
    text = " ".join((raw or "").split())
    return text[:limit]


def clamp(value: float, low: float, high: float) -> float:
    return max(low, min(high, value))


def load_json(path: Path) -> Any:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except FileNotFoundError:
        raise SystemExit(f"[inject_pptx_anchors] file not found: {path}")
    except json.JSONDecodeError as exc:
        raise SystemExit(f"[inject_pptx_anchors] invalid JSON {path}: {exc}")


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def node_text(value: Any) -> str:
    return clean_text(str(value or ""))


def c_nv_pr(shape: Any) -> Any | None:
    element = getattr(shape, "element", None)
    if element is None:
        return None
    try:
        iterator = element.iter()
    except Exception:
        return None
    for node in iterator:
        if local_name(str(getattr(node, "tag", ""))) == "cNvPr":
            return node
    return None


def shape_metadata(shape: Any) -> str:
    values: list[str] = []
    try:
        values.append(str(shape.name))
    except Exception:
        pass
    node = c_nv_pr(shape)
    if node is not None:
        for key in ("name", "title", "descr"):
            values.append(str(node.get(key) or ""))
    try:
        if getattr(shape, "has_text_frame", False):
            values.append(str(shape.text))
    except Exception:
        pass
    return " ".join(values)


def existing_anchor_ids(prs: Any) -> set[str]:
    anchors: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            anchors.update(ANCHOR_RE.findall(shape_metadata(shape)))
    return anchors


def set_shape_metadata(shape: Any, anchor_id: str, text: str) -> None:
    try:
        shape.name = anchor_id
    except Exception:
        pass
    node = c_nv_pr(shape)
    if node is not None:
        node.set("name", anchor_id)
        node.set("title", anchor_id)
        node.set("descr", clean_text(f"{anchor_id} {text}", limit=320))


def make_shape_invisible(shape: Any) -> None:
    """Hide a rectangle in normal rendering while retaining PPTX geometry."""
    try:
        from pptx.dml.color import RGBColor  # type: ignore

        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        try:
            shape.fill.transparency = 100
        except Exception:
            pass
        try:
            shape.line.fill.background()
        except Exception:
            pass
    except Exception:
        pass

    # Make transparency explicit for Office/LibreOffice readers that ignore the
    # high-level python-pptx transparency helper.
    try:
        from pptx.oxml.ns import qn  # type: ignore

        sp_pr = shape._element.spPr  # noqa: SLF001 - python-pptx exposes XML here.
        solid = sp_pr.find(qn("a:solidFill"))
        if solid is None:
            solid = sp_pr._new_solidFill()  # noqa: SLF001
            sp_pr.insert(0, solid)
        srgb = solid.find(qn("a:srgbClr"))
        if srgb is None:
            srgb = solid._new_srgbClr("FFFFFF")  # noqa: SLF001
            solid.append(srgb)
        for alpha in list(srgb.findall(qn("a:alpha"))):
            srgb.remove(alpha)
        alpha = srgb.makeelement(qn("a:alpha"), {"val": "0"})
        srgb.append(alpha)

        ln = sp_pr.find(qn("a:ln"))
        if ln is not None:
            sp_pr.remove(ln)
    except Exception:
        pass


def cue_entries(cue_plan: dict[str, Any], min_box: float) -> list[dict[str, Any]]:
    entries: list[dict[str, Any]] = []
    for slide in cue_plan.get("slides") or []:
        if not isinstance(slide, dict):
            continue
        try:
            slide_index = int(slide.get("index"))
        except (TypeError, ValueError):
            continue
        for chunk in slide.get("chunks") or []:
            if not isinstance(chunk, dict):
                continue
            anchor_id = node_text(chunk.get("anchor_id"))
            if not anchor_id:
                continue
            box = chunk.get("region_box")
            if isinstance(box, list) and len(box) == 4:
                try:
                    x, y, w, h = [float(v) for v in box]
                except (TypeError, ValueError):
                    continue
            else:
                point = chunk.get("point")
                if not (isinstance(point, list) and len(point) == 2):
                    continue
                try:
                    px, py = [float(v) for v in point]
                except (TypeError, ValueError):
                    continue
                w = h = min_box
                x = px - w / 2.0
                y = py - h / 2.0
            w = max(w, min_box)
            h = max(h, min_box)
            x = clamp(x, 0.0, 1.0 - w)
            y = clamp(y, 0.0, 1.0 - h)
            entries.append({
                "slide_index": slide_index,
                "anchor_id": anchor_id,
                "chunk_index": chunk.get("chunk_index"),
                "chunk_id": chunk.get("chunk_id"),
                "text": clean_text(chunk.get("text") or ""),
                "box": [x, y, w, h],
                "target_source": chunk.get("target_source"),
            })
    return entries


def save_presentation(prs: Any, output: Path, source: Path) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    if output.resolve() == source.resolve():
        with tempfile.NamedTemporaryFile(prefix="pptx_anchor_", suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        try:
            prs.save(str(tmp_path))
            shutil.move(str(tmp_path), str(output))
        finally:
            if tmp_path.exists():
                tmp_path.unlink()
    else:
        prs.save(str(output))


def inject(pptx_path: Path, cue_plan_path: Path, output_path: Path, report_path: Path | None, min_box: float) -> dict[str, Any]:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
        from pptx.util import Emu  # type: ignore
    except ImportError as exc:
        raise SystemExit(
            "[inject_pptx_anchors] python-pptx is not installed in this Python env. "
            "Run inside the ACL26 environment."
        ) from exc

    if not pptx_path.is_file():
        raise SystemExit(f"[inject_pptx_anchors] PPTX not found: {pptx_path}")
    cue_plan = load_json(cue_plan_path)
    if not isinstance(cue_plan, dict) or not isinstance(cue_plan.get("slides"), list):
        raise SystemExit("[inject_pptx_anchors] cue plan must contain a slides array")

    prs = Presentation(str(pptx_path))
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    existing = existing_anchor_ids(prs)
    injected: list[dict[str, Any]] = []
    skipped: list[dict[str, Any]] = []

    for entry in cue_entries(cue_plan, min_box=min_box):
        anchor_id = entry["anchor_id"]
        slide_index = entry["slide_index"]
        if anchor_id in existing:
            skipped.append({**entry, "reason": "anchor_already_present"})
            continue
        if slide_index < 1 or slide_index > len(prs.slides):
            skipped.append({**entry, "reason": "slide_index_out_of_range"})
            continue
        x, y, w, h = entry["box"]
        slide = prs.slides[slide_index - 1]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(x * slide_w)),
            Emu(int(y * slide_h)),
            Emu(int(w * slide_w)),
            Emu(int(h * slide_h)),
        )
        set_shape_metadata(shape, anchor_id, entry.get("text") or "")
        make_shape_invisible(shape)
        existing.add(anchor_id)
        injected.append(entry)

    save_presentation(prs, output_path, pptx_path)
    report = {
        "schema_version": SCHEMA_VERSION,
        "created_at": utc_now(),
        "pptx": str(pptx_path),
        "output": str(output_path),
        "cue_plan": str(cue_plan_path),
        "slide_count": len(prs.slides),
        "injected_count": len(injected),
        "skipped_count": len(skipped),
        "injected": injected,
        "skipped": skipped,
    }
    if report_path:
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description="Inject cue-plan anchors into PPTX shape metadata.")
    parser.add_argument("--pptx", type=Path, required=True, help="Input PPTX produced by ppt-master.")
    parser.add_argument("--cue-plan", type=Path, required=True, help="visual_cue_plan.json with region_box entries.")
    parser.add_argument("--out", type=Path, required=True, help="Output PPTX. May be the same path as --pptx.")
    parser.add_argument("--report", type=Path, default=None, help="Optional JSON report.")
    parser.add_argument("--min-box", type=float, default=0.015, help="Minimum normalized box width/height for point-only cues.")
    args = parser.parse_args()

    report = inject(args.pptx.resolve(), args.cue_plan.resolve(), args.out.resolve(), args.report.resolve() if args.report else None, args.min_box)
    print(
        f"[inject_pptx_anchors] injected {report['injected_count']} anchor(s), "
        f"skipped {report['skipped_count']} anchor(s) -> {args.out}"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
